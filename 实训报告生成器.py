
# -*- coding: utf-8 -*-
"""实训报告生成器

基于实训报告模板和按天整理的 TXT/PPTX/DOCX 笔记，自动生成后续课程报告正文。

常用示例：p
    python 实训报告生成器.py
    python 实训报告生成器.py 笔记1.txt 笔记2.txt 笔记3.txt --output 班级-组号-学号-姓名.docx
    python 实训报告生成器.py --screenshots 截图文件夹 --fill-cover --college 学院 --major-class 专业班级 --student-id 学号 --name 姓名
"""

from __future__ import annotations

import argparse
import ast
import hashlib
import json
import os
import random
import re
import socket
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from textwrap import dedent
from typing import Iterable
from urllib import error as url_error
from urllib import request as url_request
from urllib.parse import urlparse
from xml.sax.saxutils import escape as xml_escape

from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt


DEFAULT_TEMPLATE_CANDIDATES = [
    Path("实训报告（过程）模板.docx"),
    Path("实训报告（过程）模板(2).docx"),
    Path(r"C:\Users\SKY\Documents\xwechat_files\wxid_v6es2tk5aij412_a2db\msg\file\2026-07\实训报告（过程）模板(2).docx"),
    Path(r"C:\Users\SKY\Documents\xwechat_files\wxid_v6es2tk5aij412_a2db\msg\file\2026-07\实训报告（过程）模板.docx"),
]

DEFAULT_SOURCE_CANDIDATES = [
    Path(r"C:\Users\SKY\Desktop\学习笔记7-5-01.txt"),
    Path(r"C:\Users\SKY\Desktop\学习笔记7-6-01.txt"),
]

TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "gb18030", "gbk", "cp936")
BODY_SIZE = 10.5
BODY_FIRST_LINE_INDENT = Pt(BODY_SIZE * 2)
CN_NUMS = "零一二三四五六七八九十"
AUTO_IMAGE_DIR = Path(__file__).resolve().parent / ".report_web_cache" / "run_images"
PYCHARM_COMMAND_LINE = r"D:\PROJECT\.venv\Scripts\python.exe D:\PROJECT\lianxi.py"
NO_TASK_MESSAGE = "这次没有任务和练习，可能老师上课布置的作业，回忆回忆，自己添加哦！"
SAFE_IMPORTS = {"random", "math", "statistics", "re"}
DANGEROUS_NAMES = {
    "__import__",
    "compile",
    "ctypes",
    "delattr",
    "eval",
    "exec",
    "exit",
    "getattr",
    "globals",
    "help",
    "locals",
    "os",
    "pathlib",
    "quit",
    "setattr",
    "shutil",
    "socket",
    "subprocess",
    "sys",
    "vars",
}
SAMPLE_NAMES = ["LinQiao", "ChenYu", "ZhaoNing", "XuRan", "LuoMing", "HeYue", "TangXin", "QinMo", "SunYi", "YeFan"]
SAMPLE_CHINESE_NAMES = ["林乔", "陈宇", "赵宁", "徐然", "罗明", "何月", "唐欣", "秦墨", "孙一", "叶凡"]
SAMPLE_CLASSES = ["23060101", "23060102", "计科1", "软件2", "大数据1"]
SAMPLE_GROUPS = ["1", "2", "3", "4", "A", "B"]
SAMPLE_PRODUCTS = ["笔记本", "资料册", "中性笔", "文件夹", "U盘", "练习本"]
SAMPLE_PROJECTS = ["Python课堂练习", "暑期实训", "数据处理基础", "程序设计练习", "课程记录"]
SAMPLE_TEXT_LINES = [
    "本次练习完成了输入、处理和输出的基本流程。",
    "课堂中重点观察了代码运行结果和变量变化。",
    "通过调试可以更清楚地理解程序执行顺序。",
]
EXPLICIT_TASK_PATTERN = r"(?:实训任务|课堂任务|课堂练习|上机练习|课后练习|练习|任务)"
CODE_ACTION_KEYWORDS = (
    "代码",
    "程序",
    "编写",
    "实现",
    "运行",
    "调试",
    "输出",
    "输入",
    "绘制",
    "生成",
    "调用",
    "定义",
    "遍历",
    "读取",
    "写入",
)
PYTHON_CODE_KEYWORDS = (
    "print",
    "input",
    "for",
    "while",
    "if",
    "elif",
    "else",
    "def",
    "return",
    "import",
    "range",
    "type",
    "open",
    "read",
    "write",
    "upper",
    "lower",
    "split",
    "join",
    "find",
    "count",
    "replace",
    "strip",
    "append",
    "sort",
    "len",
    "plt",
    "matplotlib",
    "pandas",
    "numpy",
)
PARALLEL_CODE_TOPIC_ALIASES = [
    ("柱状图", ("柱状图", "条形图", "bar")),
    ("折线图", ("折线图", "曲线图", "line")),
    ("饼图", ("饼图", "圆饼图", "pie")),
    ("散点图", ("散点图", "scatter")),
    ("直方图", ("直方图", "hist")),
    ("箱线图", ("箱线图", "箱型图", "boxplot")),
    ("雷达图", ("雷达图", "radar")),
]
GENERATION_WARNINGS: list[str] = []
AI_USAGE_EVENTS: list[str] = []


@dataclass
class TaskItem:
    title: str
    requirement: str
    code: str
    caption: str


@dataclass
class DayReport:
    source: Path
    day_index: int
    title: str
    goals: list[str]
    topics: list[str]
    details: list[str]
    tasks: list[TaskItem]


@dataclass
class SourceGroup:
    label: str
    paths: list[Path]
    sort_key: tuple


@dataclass
class ComparisonRow:
    project: str
    original_text: str
    generated_text: str


@dataclass
class AIConfig:
    enabled: bool = False
    base_url: str = ""
    api_key: str = ""
    model: str = ""
    temperature: float = 0.2
    timeout: int = 60
    proxy_url: str = ""


COMMON_LOCAL_PROXY_PORTS = (7890, 7897, 10809, 1080, 20171, 6152)
COMMON_LOCAL_PROXY_HOSTS = ("127.0.0.1", "localhost")


TOPIC_RULES = [
    ("实训规划与注意事项", ("实训的规划", "实训规划", "注意事项")),
    ("Python编程语言与业务应用", ("python编程语言工具", "可以完成的业务", "业务之间的关系")),
    ("数据类型书写与元素访问", ("数据类型的书写", "元素的访问", "数据类型")),
    ("输入输出命令", ("输入/输出", "输入输出", "print", "input")),
    ("字符串大小写转换", ("upper()", "lower()", "大写", "小写")),
    ("字符串查找与统计", ("find(", "count(", "查找", "统计")),
    ("字符串替换", ("replace(", "替换")),
    ("字符串去空", ("strip(", "去掉字符串前后", "前后空格")),
    ("字符串拆分与组合", ("split(", "join(", "拆分", "连接成字符串")),
    ("字符串类型测试", ("isalnum", "isalpha", "isdigit", "isspace", "islower", "isupper", "istitle")),
    ("字符编码转换", ("ord(", "chr(", "编码转换")),
    ("运算符和表达式", ("运算符", "表达式")),
    ("选择结构", ("选择结构", "if", "elif", "else")),
    ("循环结构", ("循环结构", "for", "while", "九九乘法表", "猜数字")),
    ("break与continue", ("break", "continue")),
    ("模块导入", ("import", "from", "random", "导入python功能库")),
    ("函数定义与调用", ("函数定义", "def", "函数的调用")),
    ("函数返回值", ("return", "返回值")),
    ("函数参数传递", ("形参", "实参", "参数的传递")),
    ("变量作用域", ("全局变量", "局部变量", "变量的作用域")),
    ("文件读写命令", ("with", "open", "read()", "write()", "文件读写")),
    ("文件路径", ("绝对路径", "相对路径")),
]

DETAIL_BANK = {
    "实训规划与注意事项": "实训开始阶段需要明确课程安排、学习纪律、阶段任务、报告提交与答辩要求，为后续项目实践建立清晰的学习路线。",
    "Python编程语言与业务应用": "Python是本次实训的主要编程工具，可用于数据处理、自动化办公、数据可视化、机器学习、深度学习及项目原型开发。",
    "数据类型书写与元素访问": "Python常见数据类型包括字符串、数值、布尔值、列表、字典、元组和集合，序列类型可通过下标和切片访问元素，字典可通过键访问值。",
    "输入输出命令": "print()用于在运行结果栏输出信息，input()用于获取键盘输入，二者可结合f-string完成交互式程序设计。",
    "字符串大小写转换": "upper()和lower()用于字符串大小写转换，适合处理用户名、编码、英文文本统一格式等场景。",
    "字符串查找与统计": "find()可返回指定子串首次出现的位置，count()可统计子串出现次数，能够完成文本检索与简单统计分析。",
    "字符串替换": "replace()可把字符串中的指定内容替换为新内容，可用于文本纠错、统一术语和批量清洗字符。",
    "字符串去空": "strip()用于去除字符串首尾空白字符或指定字符，常用于清理用户输入、文件读取内容和网络文本。",
    "字符串拆分与组合": "split()可按分隔符把字符串拆分为列表，join()可把序列元素重新组合为字符串，是文本结构化处理的基础方法。",
    "字符串类型测试": "isalnum()、isalpha()、isdigit()等方法返回布尔值，可用于判断字符串是否满足用户名、编号或纯数字等格式要求。",
    "字符编码转换": "ord()可以把字符转换为编码值，chr()可以把编码值转换为字符，帮助理解字符在计算机中的存储表示。",
    "运算符和表达式": "表达式由数据、变量和运算符组成，常见运算符包括算术运算符、关系运算符和逻辑运算符，是程序判断与计算的基础。",
    "选择结构": "选择结构通过if、elif、else按照条件执行不同语句，可实现费用判断、成绩分类、权限判断等分支逻辑。",
    "循环结构": "循环结构通过for或while重复执行代码，双层循环可控制行列输出，适合九九乘法表、累加、遍历列表等任务。",
    "break与continue": "break用于提前结束循环，continue用于跳过本轮循环后继续下一轮，可实现更加灵活的循环控制。",
    "模块导入": "import和from...import用于导入标准库或第三方库，导入后即可调用库中的函数扩展程序功能。",
    "函数定义与调用": "函数通过def定义，将重复使用的功能封装为独立代码块，调用函数时传入实参即可执行对应功能。",
    "函数返回值": "return用于把函数处理结果返回给调用位置，若函数没有显式return，则默认返回None。",
    "函数参数传递": "参数传递的本质是把实参提供给形参，使函数能够处理不同输入，提高代码复用性。",
    "变量作用域": "全局变量定义在函数外，局部变量定义在函数内，理解作用域有助于避免变量覆盖和数据混乱。",
    "文件读写命令": "with open()可安全打开文件，read()、readline()、readlines()用于读取，write()用于写入文本内容。",
    "文件路径": "绝对路径从磁盘根目录开始定位文件，相对路径以当前程序运行目录为参照，二者都会影响文件能否正确打开。",
}

GOAL_BANK = {
    "实训规划与注意事项": "了解暑期实训的整体安排、学习纪律、报告提交和项目答辩要求。",
    "Python编程语言与业务应用": "认识Python编程语言在数据分析、人工智能和项目实践中的工具作用。",
    "数据类型书写与元素访问": "掌握Python常见数据类型的书写方式，并能够使用下标、切片或键访问元素。",
    "输入输出命令": "能够使用print()、input()和格式化输出完成基础交互程序。",
    "字符串大小写转换": "掌握字符串大小写转换方法，并能完成用户名等文本格式处理。",
    "字符串查找与统计": "能够使用find()和count()完成文本查找与次数统计。",
    "字符串替换": "能够使用replace()对文本中的指定字符或词语进行批量替换。",
    "字符串去空": "能够使用strip()清理字符串首尾空白字符。",
    "字符串拆分与组合": "能够使用split()和join()完成字符串与列表之间的转换。",
    "字符串类型测试": "能够使用字符串测试函数判断文本是否满足指定格式。",
    "字符编码转换": "理解ord()与chr()在字符和编码之间转换的作用。",
    "运算符和表达式": "理解运算符和表达式的基本构成，并能完成简单计算与判断。",
    "选择结构": "掌握if、elif、else选择结构的基本写法。",
    "循环结构": "掌握for、while循环结构，并能够编写九九乘法表和猜数字等程序。",
    "break与continue": "理解break和continue对循环执行流程的控制作用。",
    "模块导入": "掌握import和from...import导入模块的基本方法。",
    "函数定义与调用": "能够使用def定义函数，并通过实参调用函数解决重复性问题。",
    "函数返回值": "理解return返回值的作用，并能够获取函数处理结果。",
    "函数参数传递": "理解形参与实参的关系，能够编写带参数的函数。",
    "变量作用域": "区分全局变量和局部变量，避免函数内外变量使用错误。",
    "文件读写命令": "掌握with open()文件读写结构，并能够完成文本文件的创建、写入和读取。",
    "文件路径": "理解绝对路径和相对路径的区别，能够正确定位程序运行所需文件。",
}


def main() -> int:
    args = parse_args()
    try:
        template_path = resolve_template(args.template)
        source_paths = resolve_sources(args.sources)
        output_path = resolve_output_path(args)

        doc = Document(str(template_path))
        if args.fill_cover:
            fill_cover(doc, args)

        clear_template_body_after_cover(doc)
        ensure_report_styles(doc)

        source_groups = group_sources_by_day(source_paths)
        ai_config = build_ai_config_from_args(args)
        clear_generation_warnings()
        variation_seed = args.variation_seed or f"{output_path.name}-{datetime.now().isoformat(timespec='microseconds')}"
        comparison_rows: list[ComparisonRow] = []
        for idx, group in enumerate(source_groups, start=1):
            raw_text = read_source_group_text(group)
            title_override = args.day_title[idx - 1] if idx <= len(args.day_title) else None
            report = build_day_report(
                source=group.paths[0],
                day_index=idx,
                raw_text=raw_text,
                title_override=title_override,
                max_tasks=args.max_tasks_per_day,
                ai_config=ai_config,
                variation_seed=f"{variation_seed}-{group.label}-{idx}",
            )
            comparison_rows.extend(build_comparison_rows(report, raw_text, group))
            if idx > 1 and args.day_page_break:
                doc.add_page_break()
            append_day_report(
                doc,
                report,
                screenshots_dir=Path(args.screenshots) if args.screenshots else None,
                auto_result_images=not args.no_auto_result_images,
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
        comparison_path = output_path.with_name("对照表.xlsx")
        write_comparison_xlsx(comparison_path, comparison_rows)

        print(f"已生成实训报告：{output_path}")
        print(f"已生成对照表：{comparison_path}")
        print(f"使用模板：{template_path}")
        print("按天使用的资料：")
        for day_index, group in enumerate(source_groups, start=1):
            print(f"  第{day_index}天（{group.label}）：")
            for source in group.paths:
                print(f"    - {source}")
        for warning in get_generation_warnings():
            print(f"提示：{warning}")
        return 0
    except Exception as exc:  # noqa: BLE001 - 命令行工具需要给出清楚错误
        print(f"生成失败：{exc}", file=sys.stderr)
        return 1


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="按实训报告模板生成课程目标、课程内容、内容详情、代码及运行结果截图占位。",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    parser.add_argument("sources", nargs="*", help="按天顺序传入 TXT/PPTX/DOCX 笔记文件。为空时使用脚本内置的两个样例 TXT。")
    parser.add_argument("--template", help="实训报告模板 DOCX 路径。为空时自动寻找同目录模板或本机默认模板。")
    parser.add_argument("-o", "--output", help="输出 DOCX 路径。为空时自动生成文件名。")
    parser.add_argument("--screenshots", help="运行结果截图文件夹；自动匹配 1-1*.png、2-3*.jpg 等文件并插入。")
    parser.add_argument("--day-title", action="append", default=[], help="覆盖某一天标题，可重复传入，顺序对应资料文件。")
    parser.add_argument("--day-page-break", action="store_true", help="每天之间强制分页。")
    parser.add_argument("--max-tasks-per-day", type=int, default=8, help="每天最多写入的任务/练习数量。")
    parser.add_argument("--no-auto-result-images", action="store_true", help="关闭代码运行结果图自动生成，只保留截图占位或手动截图。")
    parser.add_argument("--variation-seed", default="", help="差异化生成种子；为空时根据输出文件名和当前时间自动生成。")
    parser.add_argument("--fill-cover", action="store_true", help="默认不动第一页；开启后只填封面学院、专业班级、学号、姓名。")
    parser.add_argument("--college", default="", help="封面学院。")
    parser.add_argument("--major-class", default="", help="封面专业班级。")
    parser.add_argument("--class-name", default="", help="报告命名用班级，例如：计科1班。")
    parser.add_argument("--group", default="", help="报告命名用组别，例如：2组。")
    parser.add_argument("--student-id", default="", help="封面学号和报告命名用学号。")
    parser.add_argument("--name", default="", help="封面姓名和报告命名用姓名。")
    parser.add_argument("--ai", action="store_true", help="启用 OpenAI 兼容接口增强课程总结与代码生成。")
    parser.add_argument("--ai-base-url", default="", help="OpenAI 兼容 API 地址，例如：https://example.com/v1。")
    parser.add_argument("--ai-api-key", default="", help="API Key；没有鉴权要求的本地/免费接口可以留空。")
    parser.add_argument("--ai-model", default="", help="模型名称。")
    parser.add_argument("--ai-temperature", type=float, default=0.2, help="AI 生成温度。")
    parser.add_argument("--ai-proxy", default="", help="可选代理地址，例如：http://127.0.0.1:7890。")
    return parser.parse_args()


def build_ai_config_from_args(args: argparse.Namespace) -> AIConfig:
    if not getattr(args, "ai", False):
        return AIConfig()
    return AIConfig(
        enabled=True,
        base_url=args.ai_base_url.strip(),
        api_key=args.ai_api_key.strip(),
        model=args.ai_model.strip(),
        temperature=args.ai_temperature,
        proxy_url=args.ai_proxy.strip(),
    )


def resolve_template(template_arg: str | None) -> Path:
    candidates = [Path(template_arg)] if template_arg else DEFAULT_TEMPLATE_CANDIDATES
    for candidate in candidates:
        if candidate and candidate.exists():
            return candidate.resolve()
    raise FileNotFoundError("没有找到实训报告模板，请使用 --template 指定模板 DOCX 路径。")


def resolve_sources(source_args: list[str]) -> list[Path]:
    if source_args:
        paths = [Path(item).resolve() for item in source_args]
    else:
        paths = discover_uploaded_txt_sources()
        if not paths:
            paths = [item.resolve() for item in DEFAULT_SOURCE_CANDIDATES if item.exists()]
    if not paths:
        raise FileNotFoundError("没有找到课堂笔记资料，请把 TXT/PPTX/DOCX 路径作为参数传入。")
    missing = [str(path) for path in paths if not path.exists()]
    if missing:
        raise FileNotFoundError("以下资料文件不存在：" + "；".join(missing))
    return paths


def discover_uploaded_txt_sources() -> list[Path]:
    txt_files = [
        path.resolve()
        for pattern in ("*.txt", "*.TXT")
        for path in Path.cwd().glob(pattern)
        if path.is_file() and looks_like_note_filename(path.name)
    ]
    return sorted(txt_files, key=natural_sort_key)


def looks_like_note_filename(filename: str) -> bool:
    stem = Path(filename).stem
    return bool(re.search(r"(学习笔记)?\d{1,2}[-_]\d{1,2}(?:[-_]\d+)?$", stem))


def natural_sort_key(path: Path) -> list[object]:
    parts = re.split(r"(\d+)", path.name)
    return [int(part) if part.isdigit() else part.lower() for part in parts]


def group_sources_by_day(paths: list[Path]) -> list[SourceGroup]:
    grouped: dict[tuple, SourceGroup] = {}
    for path in sorted(paths, key=natural_sort_key):
        month_day = extract_month_day_key(path)
        if month_day:
            month, day = month_day
            key = ("date", month, day)
            label = f"{month}-{day}"
        else:
            key = ("file", path.name.lower())
            label = path.stem
        if key not in grouped:
            grouped[key] = SourceGroup(label=label, paths=[], sort_key=key)
        grouped[key].paths.append(path)
    return sorted(grouped.values(), key=lambda group: group.sort_key)


def extract_month_day_key(path: Path) -> tuple[int, int] | None:
    name = path.stem

    year_match = re.search(r"(?:^|[^\d])(?:19|20)\d{2}[-_.年](\d{1,2})[-_.月](\d{1,2})(?=[-_.日]|[^\d]|$)", name)
    if year_match:
        return int(year_match.group(1)), int(year_match.group(2))

    month_day_match = re.search(r"(?:^|[^\d])(\d{1,2})[-_.月](\d{1,2})(?=[-_.日]|[^\d]|$)", name)
    if month_day_match:
        return int(month_day_match.group(1)), int(month_day_match.group(2))
    return None


def read_source_group_text(group: SourceGroup) -> str:
    blocks: list[str] = []
    for path in group.paths:
        text = read_source_text(path).strip()
        blocks.append(f"【资料：{path.name}】\n{text}")
    return "\n\n".join(blocks)


def resolve_output_path(args: argparse.Namespace) -> Path:
    if args.output:
        return unique_path(Path(args.output).resolve())

    if args.class_name and args.group and args.student_id and args.name:
        group = args.group if args.group.endswith("组") else f"{args.group}组"
        filename = f"{args.class_name}-{group}-{args.student_id}-{args.name}.docx"
    else:
        filename = f"实训报告_{datetime.now():%Y%m%d_%H%M%S}.docx"
    return unique_path((Path.cwd() / filename).resolve())


def unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem, suffix = path.stem, path.suffix
    for index in range(1, 1000):
        candidate = path.with_name(f"{stem}_{index}{suffix}")
        if not candidate.exists():
            return candidate
    raise RuntimeError(f"无法为输出文件生成唯一名称：{path}")


def fill_cover(doc: Document, args: argparse.Namespace) -> None:
    values = {
        "学    院": args.college,
        "专业班级": args.major_class or args.class_name,
        "学    号": args.student_id,
        "姓    名": args.name,
    }
    for paragraph in doc.paragraphs[:30]:
        text = paragraph.text
        for label, value in values.items():
            if value and label in text:
                clear_paragraph(paragraph)
                paragraph.add_run(f"{label}： {value}")
                for run in paragraph.runs:
                    set_run_font(run, size=BODY_SIZE)


def clear_template_body_after_cover(doc: Document) -> None:
    body = doc.element.body
    remove_from = None
    paragraph_elements = {paragraph._element: paragraph for paragraph in doc.paragraphs}

    for element in list(body):
        paragraph = paragraph_elements.get(element)
        if paragraph is None:
            continue
        text = paragraph.text.strip()
        if text.startswith("一、第1天") or text.startswith("第一天"):
            remove_from = element
            break

    if remove_from is None:
        saw_cover_section = False
        for element in list(body):
            paragraph = paragraph_elements.get(element)
            if paragraph is None:
                continue
            if saw_cover_section:
                remove_from = element
                break
            if paragraph._p.pPr is not None and paragraph._p.pPr.sectPr is not None:
                saw_cover_section = True

    if remove_from is None:
        doc.add_page_break()
        return

    removing = False
    for element in list(body):
        if element is remove_from:
            removing = True
        if removing and element.tag != qn("w:sectPr"):
            body.remove(element)


def ensure_report_styles(doc: Document) -> None:
    ensure_style(doc, "实训报告_正文", size=BODY_SIZE, first_line=True)
    ensure_style(doc, "实训报告_小节标题", size=BODY_SIZE, first_line=False)
    ensure_style(doc, "实训报告_代码", size=BODY_SIZE, first_line=False, left_indent=BODY_FIRST_LINE_INDENT)
    ensure_style(doc, "实训报告_图题", size=BODY_SIZE, first_line=False, alignment=WD_ALIGN_PARAGRAPH.CENTER)
    ensure_style(doc, "实训报告_日标题", size=14, first_line=False, bold=True, alignment=WD_ALIGN_PARAGRAPH.CENTER)


def ensure_style(
    doc: Document,
    name: str,
    size: float,
    first_line: bool,
    bold: bool = False,
    alignment: WD_ALIGN_PARAGRAPH | None = None,
    left_indent=Pt(0),
) -> None:
    try:
        style = doc.styles[name]
    except KeyError:
        style = doc.styles.add_style(name, WD_STYLE_TYPE.PARAGRAPH)
    style.base_style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(size)
    style.font.bold = bold
    set_style_east_asia_font(style, "宋体")

    fmt = style.paragraph_format
    fmt.space_before = Pt(0)
    fmt.space_after = Pt(0)
    fmt.line_spacing = 1.5
    fmt.first_line_indent = BODY_FIRST_LINE_INDENT if first_line else Pt(0)
    fmt.left_indent = left_indent
    if alignment is not None:
        fmt.alignment = alignment


def set_style_east_asia_font(style, east_asia_font: str) -> None:
    r_pr = style.element.rPr
    if r_pr is None:
        r_pr = OxmlElement("w:rPr")
        style.element.append(r_pr)
    r_fonts = r_pr.rFonts
    if r_fonts is None:
        r_fonts = OxmlElement("w:rFonts")
        r_pr.append(r_fonts)
    for attr in ("w:eastAsia", "w:cs"):
        r_fonts.set(qn(attr), east_asia_font)
    for attr in ("w:ascii", "w:hAnsi"):
        r_fonts.set(qn(attr), "Times New Roman")


def read_source_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".txt":
        return read_text_file(path)
    if suffix == ".pptx":
        return read_pptx_text(path)
    if suffix == ".docx":
        return read_docx_text(path)
    raise ValueError(f"暂不支持的资料格式：{path.suffix}，请使用 TXT/PPTX/DOCX。")


def read_text_file(path: Path) -> str:
    for encoding in TEXT_ENCODINGS:
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="gb18030", errors="replace")


def read_docx_text(path: Path) -> str:
    document = Document(str(path))
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(lines)


def read_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("读取 PPTX 需要安装 python-pptx；可以先把 PPT 内容整理成 TXT 再运行。") from exc

    presentation = Presentation(str(path))
    lines: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"第{slide_index}页PPT")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def build_day_report(
    source: Path,
    day_index: int,
    raw_text: str,
    title_override: str | None,
    max_tasks: int,
    ai_config: AIConfig | None = None,
    variation_seed: str = "",
    personal_note: str = "",
) -> DayReport:
    topics = extract_topics(raw_text)
    goals = generate_goals(topics)
    details = generate_details(topics, raw_text)
    tasks = extract_tasks(raw_text)
    tasks = merge_tasks(tasks)[:max_tasks]

    title = title_override or infer_day_title(source, raw_text, topics)
    report = DayReport(source, day_index, title, goals, topics, details, tasks)
    ai_enhanced = False
    if ai_config and ai_config.enabled:
        try:
            report = enhance_day_report_with_ai(report, raw_text, max_tasks, ai_config, variation_seed, personal_note)
            ai_enhanced = True
            add_ai_usage_event(f"第{day_index}天已真实调用 API 模型：{ai_config.model}")
        except Exception as exc:  # noqa: BLE001 - AI 只是增强项，失败时应继续生成本地报告。
            add_generation_warning(f"AI 增强失败，已自动改用本地规则生成：{exc}")
    return diversify_day_report(report, raw_text, variation_seed, ai_enhanced=ai_enhanced, personal_note=personal_note)


def enhance_day_report_with_ai(
    fallback: DayReport,
    raw_text: str,
    max_tasks: int,
    ai_config: AIConfig,
    variation_seed: str = "",
    personal_note: str = "",
) -> DayReport:
    if not ai_config.base_url or not ai_config.model:
        raise ValueError("启用 AI 增强时必须填写 API 地址和模型名。")

    prompt = build_ai_prompt(fallback, raw_text, max_tasks, variation_seed, personal_note)
    content = call_openai_compatible_chat(ai_config, prompt)
    payload = extract_json_object(content)
    return merge_ai_payload_into_report(fallback, payload, max_tasks)


def test_ai_connection(ai_config: AIConfig) -> str:
    if not ai_config.base_url or not ai_config.model:
        raise ValueError("请先填写 AI API 地址和模型名。")
    prompt = '请只输出 JSON：{"ok": true, "message": "连接成功"}'
    content = call_openai_compatible_chat(ai_config, prompt)
    payload = extract_json_object(content)
    message = str(payload.get("message") or "连接成功")
    return message


def build_ai_prompt(
    fallback: DayReport,
    raw_text: str,
    max_tasks: int,
    variation_seed: str = "",
    personal_note: str = "",
) -> str:
    text = raw_text.strip()
    if len(text) > 9000:
        text = text[:9000] + "\n……（资料过长，后文已截断）"
    fallback_tasks = [
        {"title": task.title, "requirement": task.requirement, "code": task.code}
        for task in fallback.tasks[:max_tasks]
    ]
    return f"""
你是暑期实训报告助手。请根据课堂 TXT 笔记，为第{fallback.day_index}天生成实训报告内容。

必须遵守：
1. 只输出一个 JSON 对象，不要 Markdown，不要解释。
2. 中文内容用于 Word 报告，语言自然、像学生实训报告。
3. 课程目标反向概括当天学习目标。
4. 课程内容只写知识点名称，用数组列出。
5. 课程内容详情用数组列出，每条是完整句子。
6. 练习和任务提取规则：如果 TXT 里明确出现“练习”“任务”“实训任务”“课堂练习”等内容，只能提取这些明确任务；如果没有明确任务，则从 TXT 中涉及代码的语句、代码片段、代码操作步骤中提取任务；如果两类内容都没有，tasks 才能为空。
7. code 必须是纯文本 Python 代码，不要图片，不要 Markdown 代码块。
8. 最多生成 {max_tasks} 个任务；如果本地规则底稿里的 tasks 为空，返回的 tasks 必须为空数组；如果底稿 tasks 是由代码类知识点提取出来的，也必须保留对应任务，不要删空。
9. 目标、详情和任务说明要主动变换句式，不要照搬固定模板。
10. Python 代码里的变量名、示例数据和输出文本要自然变化，避免和底稿完全一致。
11. 差异化参考编号：{variation_seed or fallback.source.stem}-{fallback.day_index}。
12. 重点：必须做“内容重构”，不能只改变知识点顺序。可以从学习目的、应用场景、操作流程、调试注意点、课堂理解角度重新组织表达。
13. topics 仍然只能写知识点名称，但名称可以更具体，例如“输入输出命令”可重构为“交互式输入与格式化输出”；不要写成完整句子。
14. details 要体现新的侧重点，不能逐句复述本地规则底稿；每条建议 30-80 字。
15. 若有任务，code 要可运行，变量名、示例数据、输出文案要和底稿不同，但必须符合任务要求。
16. 保持学生实训报告口吻，不要写成 AI 总结稿、宣传稿或论文。
17. 必须根据“个性化要求”微调表达角度、示例侧重点和总结语气；如果个性化要求为空，也要从 TXT 原文中提炼差异化表达。
18. 避免与本地规则底稿使用相同开头、相同结尾、相同段落顺序；如果某一条内容和底稿高度相似，请重写后再输出。
19. 如果 TXT 中出现“生成柱状图、折线图、饼图、散点图”等并列代码类知识点，必须拆成多个独立 task，每个 task 分别给出代码。

JSON 格式：
{{
  "title": "当天标题",
  "goals": ["课程目标1", "课程目标2"],
  "topics": ["知识点1", "知识点2"],
  "details": ["内容详情1", "内容详情2"],
  "tasks": [
    {{
      "title": "任务标题",
      "requirement": "任务要求",
      "code": "Python代码",
      "caption": "图题短名称"
    }}
  ]
}}

本地规则底稿：
{json.dumps({
    "title": fallback.title,
    "goals": fallback.goals,
    "topics": fallback.topics,
    "details": fallback.details,
    "tasks": fallback_tasks,
}, ensure_ascii=False, indent=2)}

课堂 TXT 笔记：
{text}

个性化要求：
{personal_note.strip() or "无。请根据 TXT 文件名、知识点、任务类型和课堂笔记细节自行调整表达角度。"}
""".strip()


def call_openai_compatible_chat(ai_config: AIConfig, prompt: str) -> str:
    endpoint = normalize_chat_endpoint(ai_config.base_url)
    body = build_chat_body(ai_config, prompt, use_json_response_format=True)
    try:
        response_text = post_chat_request(endpoint, ai_config, body)
    except RuntimeError as exc:
        message = str(exc)
        if "response_format" not in message and "json_object" not in message and "HTTP 400" not in message and "HTTP 422" not in message:
            raise
        fallback_body = build_chat_body(ai_config, prompt, use_json_response_format=False)
        response_text = post_chat_request(endpoint, ai_config, fallback_body)

    return extract_chat_content(response_text)


def call_openai_compatible_text(
    ai_config: AIConfig,
    prompt: str,
    system_prompt: str = "你是实训报告生成器的 AI 助手，请直接回答用户问题。",
) -> str:
    endpoint = normalize_chat_endpoint(ai_config.base_url)
    body = build_chat_body(
        ai_config,
        prompt,
        use_json_response_format=False,
        system_prompt=system_prompt,
        force_json_without_response_format=False,
    )
    response_text = post_chat_request(endpoint, ai_config, body)
    return extract_chat_content(response_text)


def extract_chat_content(response_text: str) -> str:
    payload = json.loads(response_text)
    choices = payload.get("choices") or []
    if not choices:
        raise RuntimeError("AI 接口没有返回 choices。")
    message = choices[0].get("message") or {}
    content = message.get("content") or choices[0].get("text") or ""
    if isinstance(content, list):
        content = "".join(str(item.get("text") or item.get("content") or "") if isinstance(item, dict) else str(item) for item in content)
    if not str(content).strip():
        raise RuntimeError("AI 接口返回内容为空。")
    return str(content)


def build_chat_body(
    ai_config: AIConfig,
    prompt: str,
    use_json_response_format: bool,
    system_prompt: str | None = None,
    force_json_without_response_format: bool = True,
) -> dict:
    if system_prompt is None:
        system_prompt = "你只输出严格 JSON。不要输出 Markdown、解释、前后缀。"
    if not use_json_response_format and force_json_without_response_format:
        system_prompt += " 即使接口不支持 response_format，你也必须只输出可被 json.loads 解析的 JSON 对象。"
    body = {
        "model": ai_config.model,
        "temperature": ai_config.temperature,
        "stream": False,
        "max_tokens": 6000,
        "messages": [
            {
                "role": "system",
                "content": system_prompt,
            },
            {"role": "user", "content": prompt},
        ],
    }
    if use_json_response_format:
        body["response_format"] = {"type": "json_object"}
    return body


def post_chat_request(endpoint: str, ai_config: AIConfig, body: dict) -> str:
    data = json.dumps(body, ensure_ascii=False).encode("utf-8")
    headers = {
        "Content-Type": "application/json",
        "User-Agent": "ReportGenerator/1.0",
        "HTTP-Referer": "http://127.0.0.1:8765",
        "X-Title": "ReportGenerator",
    }
    if ai_config.api_key:
        headers["Authorization"] = f"Bearer {ai_config.api_key}"
    req = url_request.Request(endpoint, data=data, headers=headers, method="POST")
    attempts = build_network_attempts(ai_config)
    errors: list[str] = []
    last_error: Exception | None = None

    for label, proxy in attempts:
        try:
            opener = build_url_opener(proxy)
            with opener.open(req, timeout=ai_config.timeout) as response:
                if proxy:
                    add_ai_usage_event(f"AI 网络通道：{label}（{proxy}）")
                else:
                    add_ai_usage_event("AI 网络通道：直连")
                return response.read().decode("utf-8")
        except url_error.HTTPError as exc:
            detail = exc.read().decode("utf-8", errors="replace")
            raise RuntimeError(f"AI 接口请求失败：HTTP {exc.code}，{detail[:500]}") from exc
        except url_error.URLError as exc:
            last_error = exc
            errors.append(f"{label}：{format_ai_connection_error(exc.reason, concise=True)}")
            continue

    detail = "；".join(errors) if errors else "没有可用网络通道。"
    raise RuntimeError(f"AI 接口连接失败：已尝试 {len(attempts)} 个网络通道，{detail}") from last_error


def build_url_opener(proxy_url: str = ""):
    if proxy_url:
        return url_request.build_opener(url_request.ProxyHandler({"http": proxy_url, "https": proxy_url}))
    return url_request.build_opener(url_request.ProxyHandler({}))


def build_network_attempts(ai_config: AIConfig) -> list[tuple[str, str]]:
    explicit_proxy = normalize_proxy_url(ai_config.proxy_url)
    if explicit_proxy:
        return [(f"网页自定义代理", explicit_proxy)]

    attempts: list[tuple[str, str]] = []
    seen: set[str] = set()
    for label, proxy in detect_proxy_routes():
        if proxy and proxy not in seen:
            attempts.append((label, proxy))
            seen.add(proxy)
    attempts.append(("直连", ""))
    return attempts


def normalize_proxy_url(proxy_url: str) -> str:
    proxy_url = (proxy_url or "").strip()
    if not proxy_url:
        return ""
    if "://" not in proxy_url:
        proxy_url = f"http://{proxy_url}"
    parsed = urlparse(proxy_url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname or not parsed.port:
        raise ValueError("当前工具只支持 HTTP/HTTPS 代理，请填写类似 http://127.0.0.1:7890 的代理地址。")
    return proxy_url


def detect_proxy_routes() -> list[tuple[str, str]]:
    routes: list[tuple[str, str]] = []
    seen: set[str] = set()

    for key, value in url_request.getproxies().items():
        if key.lower() in {"http", "https", "all"}:
            try:
                proxy = normalize_proxy_url(value)
            except ValueError:
                continue
            if proxy and proxy not in seen:
                routes.append((f"系统/环境代理 {key}", proxy))
                seen.add(proxy)

    for host in COMMON_LOCAL_PROXY_HOSTS:
        for port in COMMON_LOCAL_PROXY_PORTS:
            if is_tcp_port_open(host, port):
                proxy = f"http://{host}:{port}"
                if proxy not in seen:
                    routes.append((f"本地代理端口 {host}:{port}", proxy))
                    seen.add(proxy)
    return routes


def is_tcp_port_open(host: str, port: int, timeout: float = 0.25) -> bool:
    try:
        with socket.create_connection((host, port), timeout=timeout):
            return True
    except OSError:
        return False


def probe_url_connectivity(base_url: str, proxy_url: str = "", timeout: int = 8) -> dict:
    endpoint = normalize_chat_endpoint(base_url)
    req = url_request.Request(endpoint, headers={"User-Agent": "ReportGenerator/1.0"}, method="GET")
    attempts = [("网页自定义代理", normalize_proxy_url(proxy_url))] if proxy_url else build_network_attempts(AIConfig())
    results = []
    for label, proxy in attempts:
        try:
            opener = build_url_opener(proxy)
            with opener.open(req, timeout=timeout) as response:
                return {"ok": True, "label": label, "proxy": proxy, "status": response.status, "message": "网络通道可达。"}
        except url_error.HTTPError as exc:
            if exc.code in {400, 401, 403, 404, 405, 422, 429}:
                return {"ok": True, "label": label, "proxy": proxy, "status": exc.code, "message": f"网络通道可达，接口返回 HTTP {exc.code}。"}
            results.append(f"{label}: HTTP {exc.code}")
        except Exception as exc:  # noqa: BLE001 - 检测接口要返回每条通道的错误
            results.append(f"{label}: {format_ai_connection_error(getattr(exc, 'reason', exc), concise=True)}")
    return {"ok": False, "label": "", "proxy": "", "status": 0, "message": "；".join(results) or "网络不可达。"}


def normalize_chat_endpoint(base_url: str) -> str:
    base_url = base_url.strip().rstrip("/")
    if not base_url:
        raise ValueError("AI API 地址不能为空。")
    if base_url.endswith("/chat/completions"):
        return base_url
    return f"{base_url}/chat/completions"


def format_ai_connection_error(reason, concise: bool = False) -> str:
    reason_text = str(reason)
    if "10013" in reason_text:
        if concise:
            return f"{reason_text}。Windows 拒绝 socket 连接，请使用代理或放行 {sys.executable}。"
        return (
            f"{reason_text}。这是 Windows 拒绝外网 socket 连接的错误，通常不是 API Key 写错。"
            f"当前服务实际使用的 Python 是：{sys.executable}。"
            "程序会自动尝试系统代理和常见 Clash/V2Ray 端口；如果仍失败，请确认代理软件已开启，"
            "或在网页 AI 配置里填写代理地址，例如 http://127.0.0.1:7890。"
        )
    if "getaddrinfo failed" in reason_text or "Name or service not known" in reason_text:
        return f"{reason_text}。请检查 AI API 地址是否填写正确，以及网络/DNS 是否可用。"
    if "timed out" in reason_text or "timeout" in reason_text.lower():
        return f"{reason_text}。连接超时，请检查网络、代理/VPN 或 DeepSeek 服务状态。"
    return reason_text


def extract_json_object(content: str) -> dict:
    text = content.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text, flags=re.IGNORECASE)
        text = re.sub(r"\s*```$", "", text)
    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise RuntimeError("AI 返回内容不是 JSON 对象。")
    return json.loads(text[start : end + 1])


def merge_ai_payload_into_report(fallback: DayReport, payload: dict, max_tasks: int) -> DayReport:
    title = cleanup_heading(str(payload.get("title") or fallback.title)) or fallback.title
    goals = sanitize_string_list(payload.get("goals"), fallback.goals, 7)
    topics = sanitize_string_list(payload.get("topics"), fallback.topics, 24)
    details = sanitize_string_list(payload.get("details"), fallback.details, 10)

    if not fallback.tasks:
        return DayReport(fallback.source, fallback.day_index, title, goals, topics, details, [])

    tasks: list[TaskItem] = []
    for item in payload.get("tasks") or []:
        if not isinstance(item, dict):
            continue
        task_title = cleanup_heading(str(item.get("title") or "课堂练习"))
        requirement = cleanup_requirement(str(item.get("requirement") or task_title))
        code = str(item.get("code") or "").strip()
        caption = cleanup_caption(str(item.get("caption") or task_title))
        if not code:
            code = generate_code(task_title, requirement)
        tasks.append(TaskItem(task_title, requirement, code, caption))
        if len(tasks) >= max_tasks:
            break
    if not tasks:
        tasks = fallback.tasks[:max_tasks]

    return DayReport(fallback.source, fallback.day_index, title, goals, topics, details, tasks)


def sanitize_string_list(value, fallback: list[str], limit: int) -> list[str]:
    if not isinstance(value, list):
        return fallback[:limit]
    items: list[str] = []
    for item in value:
        text = str(item).strip()
        if text:
            items.append(text)
    return unique_keep_order(items)[:limit] or fallback[:limit]


def clear_generation_warnings() -> None:
    GENERATION_WARNINGS.clear()
    AI_USAGE_EVENTS.clear()


def add_generation_warning(message: str) -> None:
    if message and message not in GENERATION_WARNINGS:
        GENERATION_WARNINGS.append(message)


def get_generation_warnings() -> list[str]:
    return GENERATION_WARNINGS[:]


def add_ai_usage_event(message: str) -> None:
    if message and message not in AI_USAGE_EVENTS:
        AI_USAGE_EVENTS.append(message)


def get_ai_usage_events() -> list[str]:
    return AI_USAGE_EVENTS[:]


VARIATION_PROFILES = [
    {
        "name": "概念梳理",
        "goal_prefix": "围绕概念理解和语法辨析",
        "detail_focus": "重点放在概念边界、语法格式和适用场景的整理上",
        "task_prefix": "从概念验证角度完成",
    },
    {
        "name": "实践应用",
        "goal_prefix": "结合课堂示例和实际操作",
        "detail_focus": "更关注输入、处理、输出之间的衔接和程序运行效果",
        "task_prefix": "从实际应用角度实现",
    },
    {
        "name": "调试反思",
        "goal_prefix": "通过编写、运行和调试程序",
        "detail_focus": "侧重理解容易出错的位置、运行结果观察和代码调整过程",
        "task_prefix": "在调试验证过程中完成",
    },
    {
        "name": "归纳迁移",
        "goal_prefix": "在归纳课堂知识的基础上",
        "detail_focus": "强调知识点之间的联系，以及后续任务中的迁移使用",
        "task_prefix": "围绕知识迁移要求完成",
    },
    {
        "name": "课堂记录",
        "goal_prefix": "结合课堂笔记中的操作记录",
        "detail_focus": "侧重把老师讲解、个人理解和实际操作步骤串联起来",
        "task_prefix": "按照课堂记录中的练习要求完成",
    },
    {
        "name": "问题解决",
        "goal_prefix": "围绕具体问题的拆解和实现",
        "detail_focus": "重点说明从题目要求到程序思路再到运行验证的完整过程",
        "task_prefix": "从问题拆解和结果验证角度完成",
    },
    {
        "name": "能力提升",
        "goal_prefix": "在巩固基础语法的同时",
        "detail_focus": "突出独立编写代码、检查结果和总结方法的能力提升",
        "task_prefix": "围绕能力训练目标完成",
    },
    {
        "name": "应用扩展",
        "goal_prefix": "联系实际应用场景",
        "detail_focus": "更多关注知识点可以怎样组合成可用的小功能",
        "task_prefix": "结合应用扩展思路完成",
    },
]

TOPIC_VARIANT_BANK = {
    "实训规划与注意事项": ["实训安排与学习规范", "实训流程和注意事项", "阶段任务与报告要求"],
    "Python编程语言与业务应用": ["Python工具认知与应用方向", "Python编程环境和业务场景", "Python在数据处理中的应用"],
    "数据类型书写与元素访问": ["常用数据类型与元素访问", "数据类型表达和索引访问", "序列、映射与集合基础"],
    "输入输出命令": ["交互式输入与格式化输出", "input与print基础交互", "程序输入输出流程"],
    "字符串大小写转换": ["字符串大小写处理", "英文字符串格式统一", "upper与lower方法应用"],
    "字符串查找与统计": ["文本查找与次数统计", "find与count文本处理", "字符串检索和简单统计"],
    "字符串替换": ["字符串内容替换", "replace文本修改", "文本批量替换处理"],
    "字符串去空": ["字符串空白清理", "strip方法与输入清洗", "文本首尾字符处理"],
    "字符串拆分与组合": ["字符串拆分和拼接", "split与join结构转换", "文本列表化与重组"],
    "字符串类型测试": ["字符串格式判断", "字符类型检测方法", "文本合法性校验"],
    "字符编码转换": ["字符编码与转换", "ord和chr编码理解", "字符与编码值互转"],
    "运算符和表达式": ["表达式计算和运算符", "算术关系与逻辑表达", "程序计算基础"],
    "选择结构": ["条件判断结构", "if分支逻辑", "多分支程序流程"],
    "循环结构": ["循环流程控制", "for与while重复执行", "循环结构及嵌套应用"],
    "break与continue": ["循环跳转控制", "break和continue流程调整", "循环提前结束与跳过"],
    "模块导入": ["模块导入与库函数调用", "import语句使用", "标准库功能扩展"],
    "函数定义与调用": ["函数封装与调用", "def语句和函数复用", "函数化程序设计"],
    "函数返回值": ["return返回结果", "函数结果传递", "返回值与调用位置"],
    "函数参数传递": ["形参与实参传递", "函数参数设计", "参数输入和结果计算"],
    "变量作用域": ["局部变量与全局变量", "变量可见范围", "函数内外变量作用域"],
    "文件读写命令": ["文本文件读写", "with open文件处理", "文件创建、写入和读取"],
    "文件路径": ["绝对路径与相对路径", "文件定位方式", "程序运行目录和路径"],
}

DETAIL_VARIANT_BANK = {
    "输入输出命令": [
        "本部分把 input() 获取数据和 print() 输出结果联系起来，适合完成带提示语的交互式小程序。",
        "输入输出是课堂练习的入口和出口，写程序时要注意提示信息清楚、变量接收正确、结果表达完整。",
        "通过输入输出命令可以观察程序运行过程，也能检查变量内容是否按照预期完成传递。",
    ],
    "数据类型书写与元素访问": [
        "数据类型决定了数据的组织方式，列表、元组、字典和集合在访问元素时各有不同规则。",
        "本部分重点在于区分字符串、数值、列表、字典等对象，并掌握下标、键和值之间的访问关系。",
        "理解数据类型后，后续进行循环遍历、条件判断和函数传参时会更容易确定处理对象。",
    ],
    "循环结构": [
        "循环结构用于重复执行相似操作，for 更适合固定次数或遍历序列，while 更适合条件控制。",
        "循环练习要关注循环变量变化、循环条件是否正确，以及嵌套循环中内外层的执行顺序。",
        "通过循环可以减少重复代码，也能完成乘法表、累计求和、数据遍历等典型任务。",
    ],
    "函数定义与调用": [
        "函数把重复功能封装成独立代码块，调用时只需要提供必要参数即可得到对应处理结果。",
        "本部分强调 def 定义、参数接收、函数体执行和调用位置之间的关系。",
        "使用函数能提升程序结构清晰度，也方便在多个任务中复用同一段处理逻辑。",
    ],
    "文件读写命令": [
        "文件读写把程序运行结果保存到外部文本中，适合记录学习内容、保存数据或读取已有资料。",
        "使用 with open() 时需要关注打开模式和编码方式，避免因为路径或编码问题导致读取失败。",
        "本部分练习体现了程序和文件之间的数据交换过程，是后续项目保存结果的基础。",
    ],
    "选择结构": [
        "选择结构根据条件表达式决定执行路径，适合处理成绩判断、权限判断和输入合法性判断等问题。",
        "编写分支程序时要注意条件顺序，避免前面的条件覆盖后面更具体的判断。",
        "if、elif、else 能让程序具备基本决策能力，是交互式程序中常见的流程控制方式。",
    ],
    "字符串拆分与组合": [
        "split() 能把结构化文本拆成列表，join() 则可以把多个片段重新组织成指定格式的字符串。",
        "文本拆分和组合常用于处理姓名、电话、字段信息等内容，是简单数据清洗的基础。",
        "本部分体现了字符串和列表之间的转换关系，能帮助理解文本结构化处理过程。",
    ],
    "字符串查找与统计": [
        "find() 和 count() 可以快速定位关键词并统计出现次数，适合进行简单文本分析。",
        "查找和统计方法能帮助观察字符串内部内容，为替换、切分和清洗操作提供依据。",
        "本部分练习强调返回值含义，例如 find() 找不到时返回 -1，count() 返回匹配次数。",
    ],
}


GOAL_REWRITE_TEMPLATES = [
    "{prefix}，理解{body}，并能结合课堂示例说明其使用方法。",
    "在本次学习中，{focus}，进一步掌握{body}。",
    "通过笔记整理和程序验证，形成对{body}的基本认识和应用能力。",
    "围绕当天知识点，能够把{body}转化为可运行、可检查的实践过程。",
    "从课堂讲解到课后整理，逐步明确{body}的使用条件、书写方法和常见注意点。",
    "能够根据笔记内容提炼{body}，并把相关知识点和实际练习联系起来。",
    "以小程序实现为落点，掌握{body}在输入、处理、输出流程中的作用。",
    "在理解语法规则的同时，能够结合运行结果判断{body}是否被正确使用。",
    "通过对示例代码的拆解，梳理{body}的操作步骤，并提升独立编写程序的能力。",
    "把{body}放到实际问题中理解，能够根据任务要求选择合适的实现方式。",
]

DETAIL_REWRITE_TEMPLATES = [
    "{body}，这一部分在报告中从{focus}展开说明。",
    "课堂内容不是孤立记忆，{body}，需要结合示例运行结果一起理解。",
    "从学习过程看，{body}，后续写代码时要注意步骤衔接和结果验证。",
    "{body}，这一知识点可以和当天其他内容组合起来完成小型程序练习。",
    "围绕该知识点，课堂上既关注语法写法，也关注实际运行后的输出表现：{body}。",
    "在整理报告时，可将其理解为“概念说明、代码实现、结果观察”三步：{body}。",
    "{body}，这类内容适合通过多写几组示例来比较不同写法之间的差别。",
    "本部分的学习重点不只是记住命令，还要能说明命令为什么这样写、运行后会得到什么结果：{body}。",
    "从任务完成角度看，{body}，它能帮助后续练习更快定位输入、变量和输出之间的关系。",
    "{body}，报告中需要把知识点和课堂练习连接起来，避免只停留在概念罗列。",
]

TASK_REWRITE_TEMPLATES = [
    "{prefix}：{body}。",
    "根据课堂任务要求，{body}，并观察程序输出是否正确。",
    "围绕本题要求编写 Python 程序，实现{body}。",
    "在 PyCharm 中完成代码编写和运行验证，任务内容为：{body}。",
    "按照题目描述先确定输入和处理步骤，再用 Python 完成{body}。",
    "本练习需要把课堂知识点落实到代码中，完成{body}并保留运行结果。",
    "根据笔记中的练习要求，设计变量、语句和输出格式，完成{body}。",
    "先分析任务需要哪些数据，再编写程序实现{body}，最后检查输出内容。",
    "结合当天知识点完成一个可运行示例，具体要求是：{body}。",
    "将题目要求拆分为数据准备、逻辑处理和结果展示三个部分，完成{body}。",
]

SUMMARY_LINK_TEMPLATES = [
    "结合原始笔记中的“{note}”，本次报告从{focus}进行整理。",
    "原 TXT 中提到“{note}”，报告将其放入当天知识点体系中重新归纳。",
    "围绕笔记记录的“{note}”，本次内容更强调理解、练习和运行结果之间的对应关系。",
    "从“{note}”这一记录可以看出，当天学习重点需要和程序实践结合起来说明。",
]

CONNECTOR_POOL = ["同时", "另外", "进一步看", "从课堂记录看", "在实际练习中", "整理笔记时", "结合运行结果", "换一个角度看"]

MODULE_ORDER_STRATEGIES = ["keep", "rotate", "shuffle-middle", "reverse-pairs"]

PERSONAL_NOTE_TEMPLATES = [
    "结合个人补充要求“{note}”，本次报告在表述时更偏向{focus}。",
    "根据用户补充的“{note}”，相关内容会适当突出课堂理解和个人整理过程。",
    "围绕“{note}”这一补充方向，报告会把知识点、练习和运行结果联系得更紧一些。",
]

SECTION_CLOSING_TEMPLATES = [
    "整体来看，这部分内容更适合通过“知识点整理加代码验证”的方式掌握。",
    "因此，本部分不只记录概念，还需要把课堂操作过程和结果截图结合起来说明。",
    "从实训报告角度看，重点在于说明自己如何理解、如何实现以及如何检查结果。",
    "后续练习中，可以继续围绕这些知识点做小规模迁移应用。",
]

GOAL_BOUNDARY_TEMPLATES = [
    "本次实训先从课堂笔记中的关键内容入手，重点掌握{body}。",
    "围绕当天学习内容，首要目标是把{body}和具体练习过程对应起来。",
    "从个人整理角度看，本次需要先明确{body}，再结合程序运行结果进行验证。",
    "在本次课程记录中，{body}是后续理解代码示例的重要基础。",
]

DETAIL_BOUNDARY_TEMPLATES = [
    "作为内容详情的展开，{body}",
    "换一个角度理解，{body}",
    "结合课堂操作过程，{body}",
    "从问题分析角度看，{body}",
]

REPETITION_THRESHOLD = 0.22
MAX_DIVERSIFY_ATTEMPTS = 8


def diversify_day_report(
    report: DayReport,
    raw_text: str,
    variation_seed: str = "",
    ai_enhanced: bool = False,
    personal_note: str = "",
) -> DayReport:
    seed_text = variation_seed or f"{report.source.name}-{report.day_index}-{datetime.now().isoformat(timespec='microseconds')}"
    best_report = report
    best_score = 1.0
    for attempt in range(MAX_DIVERSIFY_ATTEMPTS):
        candidate = build_diversified_report_once(
            report=report,
            raw_text=raw_text,
            seed_text=f"{seed_text}-{attempt}",
            ai_enhanced=ai_enhanced,
            personal_note=personal_note,
        )
        score = ngram_repetition_rate(report_text_for_repetition(candidate), n=6)
        if score < best_score:
            best_report = candidate
            best_score = score
        if score <= REPETITION_THRESHOLD and not has_repeated_openings(candidate.goals + candidate.details):
            return candidate
    return best_report


def build_diversified_report_once(
    report: DayReport,
    raw_text: str,
    seed_text: str,
    ai_enhanced: bool,
    personal_note: str,
) -> DayReport:
    digest = hashlib.sha256((seed_text + raw_text[:2000] + personal_note).encode("utf-8", errors="ignore")).hexdigest()
    rng = random.Random(int(digest[:16], 16))
    profile = VARIATION_PROFILES[rng.randrange(len(VARIATION_PROFILES))]

    if ai_enhanced:
        goals = postprocess_module_items(unique_keep_order(report.goals)[:7], rng, personal_note, profile, "goal")
        topics = reshape_item_order(unique_keep_order(report.topics)[:24], rng)
        details = postprocess_module_items(unique_keep_order(report.details)[:10], rng, personal_note, profile, "detail")
        tasks = [
            TaskItem(
                title=cleanup_heading(task.title),
                requirement=vary_task_requirement(cleanup_requirement(task.requirement), rng, profile, index),
                code=diversify_code(task.code, rng, index),
                caption=cleanup_caption(task.caption or task.title),
            )
            for index, task in enumerate(report.tasks)
        ]
        return DayReport(
            report.source,
            report.day_index,
            report.title,
            goals,
            topics,
            details,
            tasks,
        )

    topics = diversify_topics(report.topics, raw_text, rng, profile)
    goals = diversify_goals(report.goals, topics, rng, profile)
    details = diversify_details(report.details, report.topics, raw_text, rng, profile)
    goals = postprocess_module_items(goals, rng, personal_note, profile, "goal")
    topics = reshape_item_order(topics, rng)
    details = postprocess_module_items(details, rng, personal_note, profile, "detail")
    tasks = [
        TaskItem(
            title=task.title,
            requirement=vary_task_requirement(task.requirement, rng, profile, index),
            code=diversify_code(task.code, rng, index),
            caption=task.caption,
        )
        for index, task in enumerate(report.tasks)
    ]
    return DayReport(report.source, report.day_index, report.title, goals, topics, details, tasks)


def postprocess_module_items(
    items: list[str],
    rng: random.Random,
    personal_note: str,
    profile: dict,
    role: str,
) -> list[str]:
    result = reshape_item_order(unique_keep_order(items), rng)
    result = vary_module_boundaries(result, rng, role)
    result = avoid_adjacent_same_opening(result, rng, role)
    if personal_note.strip() and result:
        note = cleanup_heading(personal_note.strip())[:80]
        template = PERSONAL_NOTE_TEMPLATES[rng.randrange(len(PERSONAL_NOTE_TEMPLATES))]
        insert_at = rng.randrange(min(len(result), 3) + 1)
        result.insert(insert_at, template.format(note=note, focus=profile["detail_focus"]))
    if role == "detail" and result and rng.random() < 0.55:
        result.append(SECTION_CLOSING_TEMPLATES[rng.randrange(len(SECTION_CLOSING_TEMPLATES))])
    return unique_keep_order(result)[:10 if role == "detail" else 7]


def vary_module_boundaries(items: list[str], rng: random.Random, role: str) -> list[str]:
    if not items:
        return items
    result = items[:]
    if role == "goal":
        body = normalize_goal_body(result[0])
        template = GOAL_BOUNDARY_TEMPLATES[rng.randrange(len(GOAL_BOUNDARY_TEMPLATES))]
        result[0] = template.format(body=body)
    elif role == "detail":
        body = normalize_sentence_body(result[0].rstrip("。；; "), "detail")
        template = DETAIL_BOUNDARY_TEMPLATES[rng.randrange(len(DETAIL_BOUNDARY_TEMPLATES))]
        result[0] = template.format(body=body).rstrip("。") + "。"
    return result


def reshape_item_order(items: list[str], rng: random.Random) -> list[str]:
    if len(items) <= 2:
        return items[:]
    strategy = MODULE_ORDER_STRATEGIES[rng.randrange(len(MODULE_ORDER_STRATEGIES))]
    result = items[:]
    if strategy == "rotate":
        offset = rng.randrange(1, len(result))
        result = result[offset:] + result[:offset]
    elif strategy == "shuffle-middle":
        middle = result[1:-1]
        rng.shuffle(middle)
        result = [result[0], *middle, result[-1]]
    elif strategy == "reverse-pairs":
        for index in range(0, len(result) - 1, 2):
            if rng.random() < 0.65:
                result[index], result[index + 1] = result[index + 1], result[index]
    return result


def avoid_adjacent_same_opening(items: list[str], rng: random.Random, role: str) -> list[str]:
    if len(items) < 2:
        return items
    result = items[:]
    for index in range(1, len(result)):
        if sentence_opening(result[index]) == sentence_opening(result[index - 1]):
            body = normalize_sentence_body(result[index].rstrip("。；; "), role)
            connector = CONNECTOR_POOL[(rng.randrange(len(CONNECTOR_POOL)) + index) % len(CONNECTOR_POOL)]
            result[index] = f"{connector}，{body}。"
    return result


def sentence_opening(text: str, width: int = 6) -> str:
    clean = re.sub(r"^[（(]\d+[)）]\s*", "", text.strip())
    clean = re.sub(r"^[一二三四五六七八九十]+[、.．]\s*", "", clean)
    return clean[:width]


def has_repeated_openings(items: list[str]) -> bool:
    openings = [sentence_opening(item) for item in items if item.strip()]
    return len(openings) != len(set(openings))


def report_text_for_repetition(report: DayReport) -> str:
    parts = [report.title, *report.goals, *report.topics, *report.details]
    for task in report.tasks:
        parts.extend([task.title, task.requirement, task.code])
    return "\n".join(parts)


def ngram_repetition_rate(text: str, n: int = 6) -> float:
    clean = re.sub(r"\s+", "", text)
    if len(clean) <= n:
        return 0.0
    grams = [clean[index : index + n] for index in range(len(clean) - n + 1)]
    if not grams:
        return 0.0
    repeated = len(grams) - len(set(grams))
    return repeated / len(grams)


def diversify_topics(topics: list[str], raw_text: str, rng: random.Random, profile: dict) -> list[str]:
    varied: list[str] = []
    for index, topic in enumerate(topics):
        variants = TOPIC_VARIANT_BANK.get(topic, [topic])
        varied.append(variants[(rng.randrange(len(variants)) + index) % len(variants)])

    extra_topics = infer_contextual_topics(raw_text, rng, profile)
    insert_at = rng.randrange(len(varied) + 1) if varied else 0
    varied[insert_at:insert_at] = extra_topics[:2]
    result = unique_keep_order(varied)
    return result[:24] or topics


def infer_contextual_topics(raw_text: str, rng: random.Random, profile: dict) -> list[str]:
    lowered = raw_text.lower()
    candidates: list[str] = []
    if "pycharm" in lowered or "运行" in raw_text or "print" in lowered:
        candidates.extend(["程序运行与结果观察", "代码调试和输出验证"])
    if "input" in lowered or "输入" in raw_text:
        candidates.extend(["数据输入和变量接收", "交互提示语设计"])
    if "列表" in raw_text or "字典" in raw_text or "元组" in raw_text or "集合" in raw_text:
        candidates.extend(["容器类型对比", "数据组织方式"])
    if "函数" in raw_text or "def" in lowered:
        candidates.extend(["函数封装思路", "功能模块化表达"])
    if "文件" in raw_text or "open" in lowered:
        candidates.extend(["文件数据保存", "读写流程验证"])
    if "调试" in raw_text or profile["name"] == "调试反思":
        candidates.extend(["错误排查和运行验证", "代码调整过程"])
    rng.shuffle(candidates)
    return unique_keep_order(candidates)


def diversify_goals(goals: list[str], topics: list[str], rng: random.Random, profile: dict) -> list[str]:
    source = goals[:] or ["梳理课堂知识点，理解相关语法并完成必要的实践记录。"]
    result: list[str] = []
    for index, goal in enumerate(source[:5]):
        body = normalize_goal_body(goal)
        template = GOAL_REWRITE_TEMPLATES[(rng.randrange(len(GOAL_REWRITE_TEMPLATES)) + index) % len(GOAL_REWRITE_TEMPLATES)]
        result.append(template.format(prefix=profile["goal_prefix"], focus=profile["detail_focus"], body=body))

    if topics:
        topic_sample = "、".join(topics[: min(3, len(topics))])
        closing_templates = [
            "能够围绕{topics}等内容完成归纳整理，并在后续练习中灵活迁移使用。",
            "能够把{topics}等知识点放到同一学习脉络中理解，形成较完整的课堂记录。",
            "能够根据{topics}等内容整理出学习重点，并用代码运行结果进行辅助说明。",
            "能够从{topics}等知识点出发，分析其在课堂练习和实训报告中的呈现方式。",
        ]
        closing = closing_templates[rng.randrange(len(closing_templates))]
        result.append(closing.format(topics=topic_sample))
    return unique_keep_order(result)[:7]


def normalize_goal_body(goal: str) -> str:
    body = normalize_sentence_body(goal.strip().rstrip("。；; "), "goal")
    body = re.sub(r"^(能够|可以|能|掌握|理解|认识|了解)", "", body).strip("，,：: ")
    return body or goal.strip().rstrip("。；; ")


def diversify_details(
    details: list[str],
    original_topics: list[str],
    raw_text: str,
    rng: random.Random,
    profile: dict,
) -> list[str]:
    result: list[str] = []
    topic_pool = original_topics[:] or extract_topics(raw_text)

    for index, topic in enumerate(topic_pool[:5]):
        variants = DETAIL_VARIANT_BANK.get(topic)
        if variants:
            result.append(variants[(rng.randrange(len(variants)) + index) % len(variants)])

    for index, detail in enumerate(details):
        result.append(rewrite_detail_sentence(detail, rng, profile, index))

    note_lines = extract_original_knowledge(raw_text, limit=360).splitlines()
    if note_lines:
        selected = note_lines[rng.randrange(len(note_lines))]
        if 8 <= len(selected) <= 80:
            template = SUMMARY_LINK_TEMPLATES[rng.randrange(len(SUMMARY_LINK_TEMPLATES))]
            result.append(template.format(note=cleanup_heading(selected), focus=profile["detail_focus"]))

    return unique_keep_order(result)[:10] or details


def rewrite_detail_sentence(text: str, rng: random.Random, profile: dict, index: int) -> str:
    body = normalize_sentence_body(text.strip().rstrip("。；; "), "detail")
    template = DETAIL_REWRITE_TEMPLATES[(rng.randrange(len(DETAIL_REWRITE_TEMPLATES)) + index) % len(DETAIL_REWRITE_TEMPLATES)]
    return template.format(body=body, focus=profile["detail_focus"])


def vary_task_requirement(text: str, rng: random.Random, profile: dict, index: int) -> str:
    body = normalize_sentence_body(text.strip().rstrip("。；; "), "requirement")
    template = TASK_REWRITE_TEMPLATES[(rng.randrange(len(TASK_REWRITE_TEMPLATES)) + index) % len(TASK_REWRITE_TEMPLATES)]
    return template.format(prefix=profile["task_prefix"], body=body)


def vary_sentence(text: str, rng: random.Random, role: str, index: int) -> str:
    body = text.strip().rstrip("。；; ")
    if not body:
        return text
    templates = {
        "goal": [
            "通过本次实训，{body}。",
            "结合课堂示例，{body}。",
            "围绕当天内容，{body}。",
            "在练习与调试过程中，{body}。",
        ],
        "detail": [
            "{body}，并通过代码示例加深理解。",
            "课堂中围绕该部分进行了讲解，主要内容为：{body}。",
            "{body}，后续练习中需要注意输入、处理和输出之间的衔接。",
            "本部分主要围绕以下内容展开：{body}。",
        ],
        "requirement": [
            "本任务要求{body}。",
            "在 PyCharm 中完成该练习：{body}。",
            "围绕课堂知识点，编写程序实现：{body}。",
            "按照题目要求，完成{body}。",
        ],
    }
    choices = templates.get(role)
    if not choices:
        return body + "。"

    template = choices[(rng.randrange(len(choices)) + index) % len(choices)]
    normalized = normalize_sentence_body(body, role)
    return template.format(body=normalized)


def normalize_sentence_body(text: str, role: str) -> str:
    text = re.sub(r"^通过本次实训，", "", text)
    text = re.sub(r"^结合课堂示例，", "", text)
    text = re.sub(r"^本任务要求", "", text)
    text = re.sub(r"^按照题目要求，完成", "", text)
    if role == "requirement":
        return text.lstrip("，,：: ")
    return text


def diversify_code(code: str, rng: random.Random, index: int) -> str:
    variant = choose_structural_code_variant(code, rng, index)
    if variant:
        code = variant
    replacements = {
        "小明": ["小林", "小周", "小华", "小陈"],
        "小花": ["小雅", "小雪", "小敏", "小雨"],
        "西安": ["西安", "成都", "武汉", "南京"],
        "北京": ["北京", "杭州", "青岛", "广州"],
        "上海": ["上海", "苏州", "天津", "深圳"],
        "2026": ["2025", "2026", "2027", "2028"],
        "实训记录.txt": ["实训记录.txt", "课堂练习记录.txt", "python实训记录.txt", "学习记录.txt"],
        "今天学习了Python文件读写命令。": [
            "今天练习了Python文件读写命令。",
            "本次实训完成了文件写入与读取操作。",
            "课堂中学习了with open()文件处理方式。",
        ],
        "with open()可以自动关闭文件。": [
            "with open()结构可以帮助程序自动关闭文件。",
            "文件处理时需要注意编码和关闭操作。",
            "读取文件时要保证路径和编码设置正确。",
        ],
        "练习本": ["练习本", "笔记本", "中性笔", "资料册"],
        "Python编程语言": ["Python编程语言", "暑期实训", "数据分析基础", "课堂练习"],
        "13444444444": ["13444444444", "13600001234", "15812345678", "17788889999"],
        "23060101": ["23060101", "23060102", "计科1", "软件2"],
    }
    result = code
    for old, values in replacements.items():
        if old in result:
            choice = values[(rng.randrange(len(values)) + index) % len(values)]
            result = result.replace(old, choice)
    result = diversify_plot_titles(result, rng, index)
    result = diversify_plot_data(result, rng)
    result = diversify_code_variables(result, rng, index)
    return result


def diversify_plot_titles(code: str, rng: random.Random, index: int) -> str:
    if "plt.title" not in code:
        return code
    title_pools = {
        "柱状图": ["小组成绩对比柱状图", "课堂练习完成情况柱状图", "数据指标柱状展示"],
        "折线图": ["学习进度折线图", "实训数据变化折线图", "每日任务完成趋势图"],
        "饼图": ["知识点占比饼图", "课堂时间分配饼图", "练习内容结构饼图"],
        "散点图": ["学习时长与成绩散点图", "任务数量与完成时间散点图", "样本数据关系散点图"],
        "直方图": ["成绩分布直方图", "样本区间分布直方图", "数据频数直方图"],
        "箱线图": ["成绩分布箱线图", "样本离散情况箱线图", "数据波动箱线图"],
        "雷达图": ["能力维度雷达图", "学习能力对比雷达图", "综合指标雷达图"],
    }
    pool: list[str] = []
    for key, values in title_pools.items():
        if key in code:
            pool = values
            break
    if not pool:
        return code
    title = pool[(rng.randrange(len(pool)) + index) % len(pool)]
    return re.sub(r'plt\.title\("([^"]*)"\)', f'plt.title("{title}")', code, count=1)


def diversify_plot_data(code: str, rng: random.Random) -> str:
    if "matplotlib" not in code and "plt." not in code:
        return code

    def replace_numeric_list(match: re.Match) -> str:
        name = match.group(1)
        original = match.group(2)
        count = max(3, len(re.findall(r"-?\d+(?:\.\d+)?", original)))
        if name in {"hours", "study_hours", "practice_hours", "time_values", "task_count"}:
            numbers = list(range(1, count + 1))
        elif name in {"rates", "percent_values", "ratio_values", "share_values"}:
            numbers = [rng.randint(12, 38) for _ in range(count)]
        elif name in {"minutes", "finish_time"}:
            start = rng.randint(22, 42)
            numbers = [start + rng.randint(6, 15) * idx for idx in range(count)]
        else:
            numbers = [rng.randint(58, 96) for _ in range(count)]
        return f"{name} = [{', '.join(str(number) for number in numbers)}]"

    return re.sub(
        r"\b(values|scores|minutes|progress|counts|rates|finish_time|task_count|hours|study_hours|practice_hours|time_values|score_values|data_values|metric_values)\s*=\s*\[([^\]]+)\]",
        replace_numeric_list,
        code,
    )


def choose_structural_code_variant(code: str, rng: random.Random, index: int) -> str:
    variants: list[str] = []
    if "matplotlib.pyplot" in code and "柱状图" in code:
        variants = [
            """
            import matplotlib.pyplot as plt

            groups = ["A组", "B组", "C组", "D组"]
            values = [78, 86, 91, 83]

            plt.rcParams["font.sans-serif"] = ["SimHei"]
            plt.bar(groups, values, color=["#0A84FF", "#30C7B5", "#FF9500", "#AF52DE"])
            plt.title("小组成绩对比柱状图")
            plt.xlabel("小组")
            plt.ylabel("成绩")
            print("柱状图已生成，最高成绩：", max(values))
            plt.show()
            """,
            """
            import matplotlib.pyplot as plt

            labels = ["理论", "练习", "调试", "整理"]
            minutes = [35, 55, 30, 25]

            plt.rcParams["font.sans-serif"] = ["SimHei"]
            plt.barh(labels, minutes, color="#5B8DEF")
            plt.title("课堂时间分配条形图")
            print("条形图已生成，总时长：", sum(minutes))
            plt.show()
            """,
        ]
    elif "matplotlib.pyplot" in code and "折线图" in code:
        variants = [
            """
            import matplotlib.pyplot as plt

            weeks = ["第1周", "第2周", "第3周", "第4周"]
            progress = [45, 62, 78, 91]

            plt.rcParams["font.sans-serif"] = ["SimHei"]
            plt.plot(weeks, progress, marker="o", linewidth=2)
            plt.title("实训进度折线图")
            plt.xlabel("时间")
            plt.ylabel("进度")
            print("折线图已生成，最终进度：", progress[-1])
            plt.show()
            """,
            """
            import matplotlib.pyplot as plt

            days = list(range(1, 6))
            counts = [3, 5, 4, 7, 8]

            plt.plot(days, counts, marker="s", color="#34C759")
            plt.title("每日练习数量变化")
            print("折线图已生成，平均数量：", sum(counts) / len(counts))
            plt.show()
            """,
        ]
    elif "matplotlib.pyplot" in code and ("饼图" in code or "pie(" in code):
        variants = [
            """
            import matplotlib.pyplot as plt

            labels = ["输入输出", "字符串", "函数", "文件读写"]
            rates = [20, 30, 25, 25]

            plt.rcParams["font.sans-serif"] = ["SimHei"]
            plt.pie(rates, labels=labels, autopct="%1.1f%%")
            plt.title("知识点占比饼图")
            print("饼图已生成，分类数量：", len(labels))
            plt.show()
            """,
            """
            import matplotlib.pyplot as plt

            items = ["讲解", "演示", "练习", "总结"]
            values = [30, 20, 35, 15]

            plt.rcParams["font.sans-serif"] = ["SimHei"]
            plt.pie(values, labels=items, startangle=120)
            print("饼图已生成，总占比：", sum(values))
            plt.show()
            """,
        ]
    elif "matplotlib.pyplot" in code and ("散点图" in code or "scatter(" in code):
        variants = [
            """
            import matplotlib.pyplot as plt

            hours = [1, 2, 3, 4, 5, 6]
            scores = [61, 69, 76, 84, 88, 94]

            plt.rcParams["font.sans-serif"] = ["SimHei"]
            plt.scatter(hours, scores, c=scores, cmap="viridis")
            plt.title("学习时长与成绩散点图")
            plt.xlabel("学习时长")
            plt.ylabel("成绩")
            print("散点图已生成，样本数量：", len(hours))
            plt.show()
            """,
            """
            import matplotlib.pyplot as plt

            task_count = [2, 3, 4, 5, 6]
            finish_time = [38, 45, 57, 70, 83]

            plt.scatter(task_count, finish_time, color="#FF9500")
            plt.title("任务数量与完成时间")
            print("散点图已生成，最大任务数：", max(task_count))
            plt.show()
            """,
        ]
    elif "转换为大写" in code and "转换为小写" in code:
        variants = [
            """
            raw_text = input("请输入一段英文或用户名：")

            print("原始内容：", raw_text)
            print("转换为大写：", raw_text.upper())
            print("转换为小写：", raw_text.lower())
            """,
            """
            source_text = input("请输入需要转换的英文内容：")
            upper_text = source_text.upper()
            lower_text = source_text.lower()

            print(f"大写结果：{upper_text}")
            print(f"小写结果：{lower_text}")
            """,
            """
            account = input("请输入账号或英文单词：")
            converted = {
                "upper": account.upper(),
                "lower": account.lower(),
            }

            for key, value in converted.items():
                print(key, "=>", value)
            """,
        ]
    elif "请输入用户名" in code and "upper" in code:
        variants = [
            """
            username = input("请输入用户名（包含数字、字母、下划线，长度为8）：")
            valid_chars = [char.isalnum() or char == "_" for char in username]

            if len(username) == 8 and all(valid_chars):
                print("转换后的用户名：", username.upper())
            else:
                print("用户名格式不符合要求")
            """,
            """
            import re

            account = input("请输入用户名（8位字母数字或下划线）：")
            pattern = r"^[A-Za-z0-9_]{8}$"

            if re.fullmatch(pattern, account):
                print("用户名大写结果：{}".format(account.upper()))
            else:
                print("输入的用户名不符合要求")
            """,
            """
            login_name = input("请输入用户名：")
            length_ok = len(login_name) == 8
            symbol_ok = login_name.replace("_", "").isalnum()

            print("校验结果：", "通过" if length_ok and symbol_ok else "未通过")
            if length_ok and symbol_ok:
                print("大写形式：", login_name.upper())
            """,
        ]
    elif "请输入班级" in code and "请输入组别" in code and "请输入姓名" in code:
        variants = [
            """
            class_name = input("请输入班级：")
            group = input("请输入组别：")
            name = input("请输入姓名：")

            print(f"我是{class_name}班的学生，我是第{group}组，我的姓名是：{name}")
            """,
            """
            info = {}
            info["班级"] = input("请输入班级：")
            info["组别"] = input("请输入组别：")
            info["姓名"] = input("请输入姓名：")

            print("学生信息：{}，{}，{}".format(info["班级"], info["组别"], info["姓名"]))
            """,
            """
            prompts = ["班级", "组别", "姓名"]
            values = [input(f"请输入{item}：") for item in prompts]
            class_no, team_no, student = values

            print(f"{student}来自{class_no}，所在小组为{team_no}。")
            """,
        ]
    elif "calc_total" in code or ("请输入商品名称" in code and "请输入商品单价" in code):
        variants = [
            """
            def calc_total(price, count):
                return price * count


            goods_name = input("请输入商品名称：")
            price = float(input("请输入商品单价："))
            count = int(input("请输入购买数量："))
            total_money = calc_total(price, count)

            print(f"{goods_name}的总价为：{total_money:.2f}元")
            """,
            """
            def get_amount(unit_price, quantity):
                amount = unit_price * quantity
                return amount


            product_name = input("请输入商品名称：")
            single_price = float(input("请输入商品单价："))
            buy_count = int(input("请输入购买数量："))

            print("{}共需支付{:.2f}元".format(product_name, get_amount(single_price, buy_count)))
            """,
            """
            item_name = input("请输入商品名称：")
            goods_price = float(input("请输入商品单价："))
            goods_count = int(input("请输入购买数量："))
            pay_money = goods_price * goods_count

            print("商品：", item_name)
            print("应付金额：", round(pay_money, 2), "元")
            """,
        ]
    elif "range(1, 10)" in code and ("*{i}" in code or "*{row}" in code or "*{col}" in code):
        variants = [
            """
            for row in range(1, 10):
                for col in range(1, row + 1):
                    print(f"{col}*{row}={row * col}", end="\\t")
                print()
            """,
            """
            for i in range(1, 10):
                line = []
                for j in range(1, i + 1):
                    line.append(f"{j}×{i}={i * j}")
                print("\\t".join(line))
            """,
            """
            rows = []
            for number in range(1, 10):
                rows.append("\\t".join(f"{left}*{number}={left * number}" for left in range(1, number + 1)))

            print("\\n".join(rows))
            """,
        ]
    elif "split(" in code and "join(" in code:
        variants = [
            """
            message = "姓名:小花，年龄:20，性别:女，电话:13444444444"
            parts = message.split("，")
            final_text = "+".join(parts)

            print("拆分后的列表：", parts)
            print("组合后的字符串：", final_text)
            """,
            """
            source_text = "姓名:小花，年龄:20，性别:女，电话:13444444444"
            info_items = source_text.split("，")

            for index, item in enumerate(info_items, start=1):
                print(f"第{index}项：{item}")
            print("连接结果：", " | ".join(info_items))
            """,
            """
            raw_info = "姓名:小花，年龄:20，性别:女，电话:13444444444"
            fields = raw_info.split("，")
            output_text = ",".join(fields)

            print(fields)
            print(output_text)
            """,
        ]
    elif "类型为" in code and "samples" in code:
        variants = [
            """
            samples = [
                ("字符类型", "Python编程语言"),
                ("数值类型", 2026),
                ("布尔类型", True),
                ("列表", ["西安", "北京", "上海"]),
                ("字典", {"姓名": "小明", "年龄": 20}),
                ("元组", ("暑期实训",)),
                ("集合", {1, 2, 2, 3}),
            ]

            for type_name, value in samples:
                print(type_name, "：", value, "，类型为：", type(value))
            """,
            """
            data_samples = {
                "字符串": "Python编程语言",
                "整数": 2026,
                "列表": ["西安", "北京", "上海"],
                "字典": {"姓名": "小明", "年龄": 20},
                "集合": {1, 2, 3},
            }

            for label, sample in data_samples.items():
                print(f"{label} -> {sample} -> {type(sample)}")
            """,
            """
            values = ["Python编程语言", 2026, True, ["西安", "北京"], {"姓名": "小明"}]

            for value in values:
                print("数据：", value)
                print("类型：", type(value).__name__)
            """,
        ]
    elif "文件内容如下" in code and "open(" in code:
        variants = [
            """
            file_name = "实训记录.txt"

            with open(file_name, "w", encoding="utf-8") as file:
                file.write("今天学习了Python文件读写命令。\\n")
                file.write("with open()可以自动关闭文件。\\n")

            with open(file_name, "r", encoding="utf-8") as file:
                content = file.read()

            print("文件内容如下：")
            print(content)
            """,
            """
            target_file = "课堂练习记录.txt"
            records = [
                "本次实训完成了文件写入操作。",
                "读取文件时需要设置合适的编码。",
            ]

            with open(target_file, "w", encoding="utf-8") as handle:
                handle.write("\\n".join(records))

            with open(target_file, encoding="utf-8") as handle:
                print("读取结果：")
                print(handle.read())
            """,
            """
            note_file = "python实训记录.txt"
            saved_text = "文件读写练习\\n使用with结构管理文件\\n"

            with open(note_file, "w", encoding="utf-8") as writer:
                writer.write(saved_text)

            with open(note_file, "r", encoding="utf-8") as reader:
                for line in reader:
                    print(line.strip())
            """,
        ]

    if not variants:
        return ""
    return dedent(variants[(rng.randrange(len(variants)) + index) % len(variants)]).strip()


def diversify_code_variables(code: str, rng: random.Random, index: int) -> str:
    variable_bank = {
        "username": ["account", "user_name", "login_name"],
        "account": ["login_name", "user_account", "input_text"],
        "converted": ["case_result", "format_map", "text_forms"],
        "raw_text": ["source_text", "input_text", "origin_text"],
        "upper_text": ["big_text", "upper_result", "uppercase_value"],
        "lower_text": ["small_text", "lower_result", "lowercase_value"],
        "class_name": ["class_no", "major_class", "class_id"],
        "group": ["group_no", "team_no", "group_id"],
        "name": ["student_name", "user_name", "person_name"],
        "news": ["news_text", "article", "football_news"],
        "mystr5": ["raw_text", "message", "source_text"],
        "mylist2": ["items", "info_list", "parts"],
        "result": ["joined_text", "output_text", "final_text"],
        "answer": ["target", "number_answer", "right_number"],
        "records": ["guess_records", "history", "guess_list"],
        "chance": ["round_no", "try_count", "step"],
        "guess": ["guess_num", "input_num", "user_guess"],
        "price": ["unit_price", "single_price", "goods_price"],
        "count": ["quantity", "buy_count", "goods_count"],
        "total": ["amount", "total_price", "money"],
        "goods_name": ["goods", "product_name", "item_name"],
        "total_money": ["pay_money", "final_money", "total_price"],
        "file_name": ["target_file", "record_file", "note_file"],
        "content": ["file_text", "record_text", "saved_text"],
        "samples": ["sample_values", "data_samples", "type_samples"],
        "type_name": ["label", "data_name", "sample_name"],
        "value": ["sample", "data_value", "item_value"],
        "groups": ["teams", "class_groups", "group_names"],
        "values": ["score_values", "data_values", "metric_values"],
        "scores": ["score_list", "grade_values", "result_scores"],
        "labels": ["category_names", "item_labels", "legend_labels"],
        "rates": ["percent_values", "ratio_values", "share_values"],
        "days": ["date_labels", "day_names", "time_points"],
        "weeks": ["week_labels", "stage_names", "periods"],
        "progress": ["progress_values", "finish_rates", "trend_values"],
        "hours": ["study_hours", "practice_hours", "time_values"],
    }
    result = code
    for old, variants in variable_bank.items():
        if not re.search(rf"\b{re.escape(old)}\b", result):
            continue
        new_name = variants[(rng.randrange(len(variants)) + index) % len(variants)]
        result = re.sub(rf"\b{re.escape(old)}\b", new_name, result)
    return result


def extract_topics(text: str) -> list[str]:
    lowered = text.lower()
    topics: list[str] = []
    for topic, keywords in TOPIC_RULES:
        if any(keyword.lower() in lowered for keyword in keywords):
            topics.append(topic)

    for line in clean_lines(text):
        match = re.match(r"^[一二三四五六七八九十]+[、.．]\s*(.+)$", line)
        if not match:
            continue
        heading = cleanup_heading(match.group(1))
        if heading and heading not in {"复习", "任务", "提交报告要求"}:
            topics.append(heading)
    return unique_keep_order(topics)


def generate_goals(topics: list[str]) -> list[str]:
    goals = [GOAL_BANK[topic] for topic in topics if topic in GOAL_BANK]
    if not goals:
        goals = ["根据课堂笔记梳理当天学习内容，掌握相关知识点的基本概念、常用语法和实践方法。"]
    return unique_keep_order(goals)[:7]


def generate_details(topics: list[str], text: str) -> list[str]:
    details = [DETAIL_BANK[topic] for topic in topics if topic in DETAIL_BANK]
    if len(details) < 4:
        for line in clean_lines(text):
            normalized = cleanup_heading(line)
            if 8 <= len(normalized) <= 80 and not is_assignment_line(normalized):
                details.append(f"课堂笔记中还涉及“{normalized}”相关内容，需要结合示例进一步理解其用法。")
            if len(details) >= 6:
                break
    return unique_keep_order(details)[:10]


def extract_tasks(text: str) -> list[TaskItem]:
    explicit_tasks = extract_explicit_tasks(text)
    if explicit_tasks:
        return explicit_tasks
    return extract_code_related_tasks(text)


def extract_explicit_tasks(text: str) -> list[TaskItem]:
    lines = clean_lines(text, keep_indent=True)
    tasks: list[TaskItem] = []
    current_title = ""
    current_lines: list[str] = []
    in_exercise = False

    def flush_current() -> None:
        nonlocal current_title, current_lines
        if current_title:
            requirement = " ".join(item.strip() for item in current_lines if item.strip())
            tasks.append(build_task(current_title, requirement or current_title))
        current_title = ""
        current_lines = []

    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            continue

        direct_task = re.match(rf"^{EXPLICIT_TASK_PATTERN}\s*(?:\d+|[一二三四五六七八九十]+)?\s*[：:、.．]\s*(.+)$", line)
        if direct_task and direct_task.group(1).strip():
            flush_current()
            in_exercise = True
            current_title = infer_task_title(direct_task.group(1))
            current_lines = [direct_task.group(1)]
            continue

        if re.match(rf"^{EXPLICIT_TASK_PATTERN}\s*(?:\d+|[一二三四五六七八九十]+)?\s*[:：]?$", line):
            in_exercise = True
            flush_current()
            continue
        if in_exercise and is_non_task_section_line(line):
            flush_current()
            in_exercise = False
            continue
        if in_exercise and re.match(r"^[一二三四五六七八九十]+[、.．]\s*", line):
            break
        if in_exercise:
            match = re.match(r"^\s*(?:[（(]?(?:\d+|[一二三四五六七八九十]+)[）).、．]|[-*])\s*(.+)$", raw_line)
            if match:
                flush_current()
                current_title = infer_task_title(match.group(1))
                current_lines = [match.group(1)]
            else:
                current_lines.append(line)

    flush_current()
    return tasks


def extract_code_related_tasks(text: str) -> list[TaskItem]:
    lines = clean_lines(text, keep_indent=True)
    tasks: list[TaskItem] = []

    tasks.extend(extract_python_snippet_tasks(lines))
    for raw_line in lines:
        line = raw_line.strip()
        if not line or is_assignment_line(line) or is_explicit_task_marker_line(line):
            continue
        if looks_like_python_code_line(line):
            continue

        parallel_topics = extract_parallel_code_topics(line)
        if len(parallel_topics) >= 2:
            for topic in parallel_topics:
                tasks.append(build_code_topic_task(topic, line))
            remainder = remove_parallel_topic_words(line)
            if is_code_related_line(remainder):
                tasks.append(build_code_line_task(remainder))
            continue

        if is_code_related_line(line):
            tasks.append(build_code_line_task(line))

    return merge_tasks(tasks)


def extract_python_snippet_tasks(lines: list[str]) -> list[TaskItem]:
    tasks: list[TaskItem] = []
    block: list[str] = []

    def flush_block() -> None:
        nonlocal block
        code = normalize_code_block(block)
        block = []
        if not code or not is_probable_python_snippet(code):
            return
        title = infer_task_title(code)
        requirement = "运行 TXT 中记录的 Python 代码片段，观察代码执行过程和输出结果。"
        tasks.append(TaskItem(title=title, requirement=requirement, code=code, caption=cleanup_caption(title)))

    for raw_line in lines:
        line = raw_line.rstrip()
        stripped = line.strip()
        if looks_like_python_code_line(stripped) or (block and line.startswith((" ", "\t"))):
            block.append(line)
            continue
        flush_block()
    flush_block()
    return merge_tasks(tasks)


def normalize_code_block(lines: list[str]) -> str:
    useful = [line.rstrip() for line in lines if line.strip()]
    if not useful:
        return ""
    min_indent = min((len(line) - len(line.lstrip(" "))) for line in useful if line.strip())
    return "\n".join(line[min_indent:] if len(line) >= min_indent else line for line in useful).strip()


def is_probable_python_snippet(code: str) -> bool:
    if len(code) < 4:
        return False
    try:
        ast.parse(code)
        return True
    except SyntaxError:
        return bool(re.search(r"\b(print|input|for|while|def|import|return|open)\b", code, flags=re.IGNORECASE))


def looks_like_python_code_line(line: str) -> bool:
    if not line:
        return False
    if re.match(r"^(?:print|input|for|while|if|elif|else|def|return|import|from|with|class)\b", line):
        return True
    if re.search(r"\b(?:plt|pd|np)\.", line):
        return True
    if re.search(r"\b(?:upper|lower|split|join|find|count|replace|strip|append|sort|read|write|open|range|type|len)\s*\(", line) and not re.search(r"[\u4e00-\u9fff]", line):
        return True
    if re.match(r"^[A-Za-z_]\w*\s*=", line) and not re.search(r"[，。；：、]", line):
        return True
    return False


def is_code_related_line(line: str) -> bool:
    lowered = line.lower()
    if extract_parallel_code_topics(line):
        return True
    if any(keyword in lowered for keyword in PYTHON_CODE_KEYWORDS):
        return True
    if ("图" in line and any(word in line for word in ("生成", "绘制", "画", "可视化", "展示"))):
        return True
    if any(action in line for action in CODE_ACTION_KEYWORDS) and ("Python" in line or "python" in lowered or "函数" in line or "变量" in line):
        return True
    return False


def is_explicit_task_marker_line(line: str) -> bool:
    return bool(re.match(rf"^{EXPLICIT_TASK_PATTERN}\s*(?:\d+|[一二三四五六七八九十]+)?\s*[:：]?$", line.strip()))


def is_non_task_section_line(line: str) -> bool:
    return bool(re.match(r"^(?:课程目标|课程内容|课程内容详情|知识点|课堂内容|学习内容|提交要求|报告要求)\s*[:：]", line.strip()))


def extract_parallel_code_topics(line: str) -> list[str]:
    lowered = line.lower()
    topics: list[str] = []
    for canonical, aliases in PARALLEL_CODE_TOPIC_ALIASES:
        if any(alias.lower() in lowered for alias in aliases):
            topics.append(canonical)
    return unique_keep_order(topics)


def remove_parallel_topic_words(line: str) -> str:
    result = line
    for _canonical, aliases in PARALLEL_CODE_TOPIC_ALIASES:
        for alias in aliases:
            if re.fullmatch(r"[A-Za-z_]+", alias):
                result = re.sub(rf"\b{re.escape(alias)}\b", "", result, flags=re.IGNORECASE)
            else:
                result = result.replace(alias, "")
    result = re.sub(r"[、,，/]+", "，", result)
    return cleanup_requirement(result)


def build_code_topic_task(topic: str, source_line: str) -> TaskItem:
    title = infer_task_title(topic)
    requirement = f"根据 TXT 中“{cleanup_requirement(source_line)}”这一并列知识点，单独编写代码完成{topic}生成，并观察运行结果。"
    return build_task(title, requirement)


def build_code_line_task(line: str) -> TaskItem:
    title = infer_task_title(line)
    requirement = f"根据 TXT 中涉及代码的内容“{cleanup_requirement(line)}”，编写 Python 示例并观察运行结果。"
    return build_task(title, requirement)


def suggest_tasks_from_keywords(text: str, existing: list[TaskItem]) -> list[TaskItem]:
    additions: list[TaskItem] = []
    existing_text = " ".join(task.title + task.requirement for task in existing)

    def add_once(title: str, requirement: str, *needles: str) -> None:
        haystack = text + existing_text
        if all(needle not in haystack for needle in needles):
            return
        if any(title in task.title for task in existing + additions):
            return
        additions.append(build_task(title, requirement))

    add_once("九九乘法表输出", "使用双层for循环打印九九乘法表，外层控制行，内层控制列。", "九九乘法表")
    add_once("猜数字游戏", "随机产生1到100之间的数字，提供6次猜测机会，记录每次猜测并提示大小关系。", "猜数字")
    add_once("break与continue循环控制", "编写循环程序，使用continue跳过部分数据，使用break提前结束循环。", "break", "continue")
    add_once("函数定义、调用与返回值", "定义带参数的函数，调用函数并使用return返回计算结果。", "函数定义", "函数的调用", "return")
    add_once("文件写入与读取", "使用with open()创建文本文件，写入实训记录后再次读取并输出。", "文件读写", "with", "open")
    return additions


def build_task(title: str, requirement: str) -> TaskItem:
    title = cleanup_heading(title)
    requirement = cleanup_requirement(requirement)
    code = generate_code(title, requirement)
    caption = cleanup_caption(title)
    return TaskItem(title=title, requirement=requirement, code=code, caption=caption)


def generate_code(title: str, requirement: str) -> str:
    text = f"{title} {requirement}"
    lowered = text.lower()
    if "柱状图" in text or "条形图" in text:
        return generate_chart_code("柱状图")
    if "折线图" in text or "曲线图" in text:
        return generate_chart_code("折线图")
    if "饼图" in text or "圆饼图" in text:
        return generate_chart_code("饼图")
    if "散点图" in text:
        return generate_chart_code("散点图")
    if "直方图" in text:
        return generate_chart_code("直方图")
    if "箱线图" in text or "箱型图" in text:
        return generate_chart_code("箱线图")
    if "雷达图" in text:
        return generate_chart_code("雷达图")
    if contains_all(text, ("用户名", "大写")):
        return dedent(
            """
            username = input("请输入用户名（包含数字、字母、下划线，长度为8）：")

            if len(username) == 8 and all(char.isalnum() or char == "_" for char in username):
                print("用户名：", username.upper())
            else:
                print("用户名格式不符合要求")
            """
        ).strip()
    if "upper" in lowered or "lower" in lowered or ("大写" in text and "小写" in text):
        return dedent(
            """
            raw_text = input("请输入一段英文或用户名：")

            print("原始内容：", raw_text)
            print("转换为大写：", raw_text.upper())
            print("转换为小写：", raw_text.lower())
            """
        ).strip()
    if "足球" in text:
        return dedent(
            """
            news = "在昨晚进行的一场足球比赛中，主场作战的球队以3比1战胜了客队。足球运动吸引了众多球迷观看，年轻球员的配合提升了足球进攻的效率。"

            print("原始新闻：", news)
            print("第1个'足球'出现的位置：", news.find("足球"))
            print("'足球'出现的次数：", news.count("足球"))
            print("替换后的新闻：", news.replace("的", "地"))
            """
        ).strip()
    if "find" in lowered or "count" in lowered or "查找" in text or "统计" in text:
        return dedent(
            """
            article = "Python课堂练习中，Python可以处理文本、数据和图表。"
            keyword = "Python"

            print("原始文本：", article)
            print("关键词首次出现位置：", article.find(keyword))
            print("关键词出现次数：", article.count(keyword))
            """
        ).strip()
    if "replace" in lowered or "替换" in text:
        return dedent(
            """
            sentence = "Python实训需要记录代码过程，代码运行结果也要保存。"
            new_sentence = sentence.replace("代码", "程序")

            print("替换前：", sentence)
            print("替换后：", new_sentence)
            """
        ).strip()
    if "strip" in text.lower() or "去掉字符串前后" in text or "前后的空格" in text:
        return dedent(
            """
            mystr5 = "\\n\\n\\n我们正在进行暑假实训\\n\\n\\n"

            print("原始字符串：", repr(mystr5))
            print("去除前后空白后：", repr(mystr5.strip()))
            """
        ).strip()
    if "split" in text.lower() or "join" in text.lower() or "拆分" in text or "连接成字符串" in text:
        return dedent(
            """
            mystr5 = "姓名:小花，年龄:20，性别:女，电话:13444444444"
            mylist2 = mystr5.split("，")
            result = "+".join(mylist2)

            print("拆分后的列表：", mylist2)
            print("组合后的字符串：", result)
            """
        ).strip()
    if "九九乘法表" in text:
        return dedent(
            """
            for i in range(1, 10):
                for j in range(1, i + 1):
                    print(f"{j}*{i}={i * j}", end="\\t")
                print()
            """
        ).strip()
    if "猜数字" in text:
        return dedent(
            """
            import random

            answer = random.randint(1, 100)
            records = []

            for chance in range(1, 7):
                guess = int(input(f"第{chance}次请输入1到100之间的数字："))
                records.append(guess)

                if guess == answer:
                    print("恭喜你猜对了")
                    break
                if guess > answer:
                    print("您猜大了")
                else:
                    print("您猜小了")

                print("剩余次数：", 6 - chance)
            else:
                print("机会用完，正确数字是：", answer)

            print("猜测记录：", records)
            """
        ).strip()
    if "break" in text or "continue" in text:
        return dedent(
            """
            for number in range(1, 21):
                if number % 2 == 0:
                    continue
                if number > 15:
                    break
                print(number, end=" ")
            """
        ).strip()
    if "函数" in text or "return" in text:
        return dedent(
            """
            def calc_total(price, count):
                total = price * count
                return total


            goods_name = input("请输入商品名称：")
            price = float(input("请输入商品单价："))
            count = int(input("请输入购买数量："))
            total_money = calc_total(price, count)

            print(f"{goods_name}的总价为：{total_money:.2f}元")
            """
        ).strip()
    if "文件" in text or "open" in text.lower() or "write" in text.lower() or "read" in text.lower():
        return dedent(
            """
            file_name = "实训记录.txt"

            with open(file_name, "w", encoding="utf-8") as file:
                file.write("今天学习了Python文件读写命令。\\n")
                file.write("with open()可以自动关闭文件。\\n")

            with open(file_name, "r", encoding="utf-8") as file:
                content = file.read()

            print("文件内容如下：")
            print(content)
            """
        ).strip()
    if "数据类型" in text or "type" in text.lower():
        return dedent(
            """
            samples = [
                ("字符类型", "Python编程语言"),
                ("数值类型", 2026),
                ("布尔类型", True),
                ("列表", ["西安", "北京", "上海"]),
                ("字典", {"姓名": "小明", "年龄": 20}),
                ("元组", ("暑期实训",)),
                ("集合", {1, 2, 2, 3}),
            ]

            for type_name, value in samples:
                print(type_name, "：", value, "，类型为：", type(value))
            """
        ).strip()
    if "列表" in text or "字典" in text or "元组" in text or "集合" in text:
        return dedent(
            """
            student = {"姓名": "陈宇", "班级": "23060101", "成绩": [86, 91, 88]}

            print("学生姓名：", student["姓名"])
            print("所在班级：", student["班级"])
            print("第一次成绩：", student["成绩"][0])
            print("平均成绩：", sum(student["成绩"]) / len(student["成绩"]))
            """
        ).strip()
    if "输入" in text and "输出" in text or "input" in lowered or "print" in lowered:
        return dedent(
            """
            class_name = input("请输入班级：")
            group = input("请输入组别：")
            name = input("请输入姓名：")

            print(f"我是{class_name}班的学生，我是第{group}组，我的姓名是：{name}")
            """
        ).strip()
    return dedent(
        f"""
        # {title}
        print("任务名称：{title}")
        print("任务要求：{requirement}")
        print("请根据课堂运行结果补充截图。")
        """
    ).strip()


def generate_chart_code(kind: str) -> str:
    if kind == "柱状图":
        body = """
        import matplotlib.pyplot as plt

        names = ["第一组", "第二组", "第三组", "第四组"]
        scores = [82, 91, 76, 88]

        plt.rcParams["font.sans-serif"] = ["SimHei"]
        plt.rcParams["axes.unicode_minus"] = False
        plt.bar(names, scores, color="#5B8DEF")
        plt.title("小组成绩柱状图")
        plt.xlabel("小组")
        plt.ylabel("成绩")
        print("柱状图已生成，数据项数量：", len(scores))
        plt.show()
        """
    elif kind == "折线图":
        body = """
        import matplotlib.pyplot as plt

        days = ["周一", "周二", "周三", "周四", "周五"]
        values = [68, 74, 81, 79, 90]

        plt.rcParams["font.sans-serif"] = ["SimHei"]
        plt.rcParams["axes.unicode_minus"] = False
        plt.plot(days, values, marker="o", color="#34C759")
        plt.title("学习进度折线图")
        plt.xlabel("日期")
        plt.ylabel("完成度")
        print("折线图已生成，最高值：", max(values))
        plt.show()
        """
    elif kind == "饼图":
        body = """
        import matplotlib.pyplot as plt

        labels = ["理论学习", "代码练习", "结果整理", "问题调试"]
        sizes = [25, 40, 20, 15]

        plt.rcParams["font.sans-serif"] = ["SimHei"]
        plt.pie(sizes, labels=labels, autopct="%1.1f%%", startangle=90)
        plt.title("课堂时间分配饼图")
        print("饼图已生成，分类数量：", len(labels))
        plt.show()
        """
    elif kind == "散点图":
        body = """
        import matplotlib.pyplot as plt

        study_hours = [1, 2, 3, 4, 5, 6]
        scores = [58, 66, 75, 82, 87, 93]

        plt.rcParams["font.sans-serif"] = ["SimHei"]
        plt.scatter(study_hours, scores, color="#FF9500")
        plt.title("学习时长与成绩散点图")
        plt.xlabel("学习时长")
        plt.ylabel("成绩")
        print("散点图已生成，样本数量：", len(study_hours))
        plt.show()
        """
    elif kind == "直方图":
        body = """
        import matplotlib.pyplot as plt

        scores = [62, 75, 81, 88, 90, 93, 70, 84, 79, 96]

        plt.rcParams["font.sans-serif"] = ["SimHei"]
        plt.hist(scores, bins=5, color="#AF52DE", edgecolor="white")
        plt.title("成绩分布直方图")
        print("直方图已生成，样本数量：", len(scores))
        plt.show()
        """
    elif kind == "箱线图":
        body = """
        import matplotlib.pyplot as plt

        scores = [62, 75, 81, 88, 90, 93, 70, 84, 79, 96]

        plt.rcParams["font.sans-serif"] = ["SimHei"]
        plt.boxplot(scores)
        plt.title("成绩分布箱线图")
        print("箱线图已生成，中位数附近数据已展示。")
        plt.show()
        """
    else:
        body = """
        import matplotlib.pyplot as plt
        import math

        labels = ["A", "B", "C", "D", "E"]
        values = [80, 90, 75, 88, 92]
        angles = [index / float(len(labels)) * 2 * math.pi for index in range(len(labels))]
        values.append(values[0])
        angles.append(angles[0])

        plt.rcParams["font.sans-serif"] = ["SimHei"]
        ax = plt.subplot(111, polar=True)
        ax.plot(angles, values, "o-", linewidth=2)
        ax.fill(angles, values, alpha=0.25)
        ax.set_thetagrids([angle * 180 / math.pi for angle in angles[:-1]], labels)
        plt.title("能力维度雷达图")
        print("雷达图已生成，维度数量：", len(labels))
        plt.show()
        """
    return dedent(body).strip()


def build_comparison_rows(report: DayReport, raw_text: str, group: SourceGroup | None = None) -> list[ComparisonRow]:
    day_label = f"第{report.day_index}天：{report.title}"
    source_files = "；".join(path.name for path in group.paths) if group else report.source.name
    original_knowledge = extract_original_knowledge(raw_text)
    source_header = f"源文件：{source_files}"
    original_tasks = "\n".join(format_original_task(task) for task in extract_tasks(raw_text))
    word_topics = "；".join(report.topics)
    word_details = "\n".join(report.details)
    generated_tasks = "\n\n".join(f"{task.title}\n{task.requirement}" for task in report.tasks)
    generated_code = "\n\n".join(f"{task.title}\n{task.code}" for task in report.tasks)

    rows = [
        ComparisonRow(
            project=f"{day_label}-课程目标",
            original_text=f"{source_header}\n\n根据 TXT 中出现的知识点、课堂操作与任务记录反向概括。",
            generated_text="\n".join(report.goals),
        ),
        ComparisonRow(
            project=f"{day_label}-课程内容",
            original_text=f"{source_header}\n\n{original_knowledge}",
            generated_text=word_topics,
        ),
        ComparisonRow(
            project=f"{day_label}-课程内容详情",
            original_text=original_knowledge,
            generated_text=word_details,
        ),
        ComparisonRow(
            project=f"{day_label}-任务与练习",
            original_text=original_tasks if report.tasks else "",
            generated_text=generated_tasks if report.tasks else "",
        ),
        ComparisonRow(
            project=f"{day_label}-代码",
            original_text=original_tasks if report.tasks else "",
            generated_text=generated_code if report.tasks else "",
        ),
    ]
    return rows


def extract_original_knowledge(raw_text: str, limit: int = 1200) -> str:
    lines = []
    in_task_area = False
    for line in clean_lines(raw_text):
        if re.match(r"^(?:练习|任务)\s*\d*\s*[:：]?", line):
            in_task_area = True
            continue
        if in_task_area and re.match(r"^[一二三四五六七八九十]+[、.．]\s*", line):
            in_task_area = False
        if in_task_area:
            continue
        if is_assignment_line(line):
            continue
        if 3 <= len(line) <= 180:
            lines.append(line)
    text = "\n".join(unique_keep_order(lines))
    return text[:limit] + ("……" if len(text) > limit else "")


def format_original_task(task: TaskItem) -> str:
    return f"{task.title}：{task.requirement}".strip("：")


def write_comparison_xlsx(path: Path, rows: list[ComparisonRow]) -> None:
    headers = ["项目", "原txt文档", "生成的word文档"]
    table = [headers]
    for row in rows:
        table.append(
            [
                row.project,
                row.original_text,
                row.generated_text,
            ]
        )
    if len(table) == 1:
        table.append(["", "", ""])

    path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", build_xlsx_content_types())
        archive.writestr("_rels/.rels", build_xlsx_root_rels())
        archive.writestr("xl/workbook.xml", build_xlsx_workbook())
        archive.writestr("xl/_rels/workbook.xml.rels", build_xlsx_workbook_rels())
        archive.writestr("xl/styles.xml", build_xlsx_styles())
        archive.writestr("xl/worksheets/sheet1.xml", build_xlsx_sheet(table))


def build_xlsx_content_types() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
  <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
  <Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/>
</Types>"""


def build_xlsx_root_rels() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>"""


def build_xlsx_workbook() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets>
    <sheet name="对照表" sheetId="1" r:id="rId1"/>
  </sheets>
</workbook>"""


def build_xlsx_workbook_rels() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>
</Relationships>"""


def build_xlsx_styles() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <fonts count="2">
    <font><sz val="11"/><name val="宋体"/></font>
    <font><b/><sz val="11"/><name val="宋体"/><color rgb="FFFFFFFF"/></font>
  </fonts>
  <fills count="3">
    <fill><patternFill patternType="none"/></fill>
    <fill><patternFill patternType="gray125"/></fill>
    <fill><patternFill patternType="solid"><fgColor rgb="FF2F5597"/><bgColor indexed="64"/></patternFill></fill>
  </fills>
  <borders count="2">
    <border><left/><right/><top/><bottom/><diagonal/></border>
    <border><left style="thin"/><right style="thin"/><top style="thin"/><bottom style="thin"/><diagonal/></border>
  </borders>
  <cellStyleXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0"/></cellStyleXfs>
  <cellXfs count="3">
    <xf numFmtId="0" fontId="0" fillId="0" borderId="1" xfId="0" applyAlignment="1"><alignment vertical="top" wrapText="1"/></xf>
    <xf numFmtId="0" fontId="1" fillId="2" borderId="1" xfId="0" applyAlignment="1"><alignment horizontal="center" vertical="center" wrapText="1"/></xf>
    <xf numFmtId="0" fontId="0" fillId="0" borderId="1" xfId="0" applyAlignment="1"><alignment vertical="top" wrapText="1"/></xf>
  </cellXfs>
  <cellStyles count="1"><cellStyle name="Normal" xfId="0" builtinId="0"/></cellStyles>
</styleSheet>"""


def build_xlsx_sheet(table: list[list[str]]) -> str:
    widths = [28, 62, 72]
    cols = "".join(f'<col min="{idx}" max="{idx}" width="{width}" customWidth="1"/>' for idx, width in enumerate(widths, start=1))
    row_xml: list[str] = []
    for row_index, row in enumerate(table, start=1):
        height = 24 if row_index == 1 else 90
        cells = []
        for col_index, value in enumerate(row, start=1):
            ref = f"{excel_col(col_index)}{row_index}"
            style = 1 if row_index == 1 else 2
            cells.append(f'<c r="{ref}" t="inlineStr" s="{style}"><is><t>{xml_escape(sanitize_xlsx_text(value))}</t></is></c>')
        row_xml.append(f'<row r="{row_index}" ht="{height}" customHeight="1">{"".join(cells)}</row>')
    dimension = f"A1:{excel_col(len(table[0]))}{len(table)}"
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <dimension ref="{dimension}"/>
  <sheetViews><sheetView workbookViewId="0"><pane ySplit="1" topLeftCell="A2" activePane="bottomLeft" state="frozen"/></sheetView></sheetViews>
  <cols>{cols}</cols>
  <sheetData>{''.join(row_xml)}</sheetData>
  <autoFilter ref="{dimension}"/>
</worksheet>"""


def excel_col(index: int) -> str:
    result = ""
    while index:
        index, remainder = divmod(index - 1, 26)
        result = chr(65 + remainder) + result
    return result


def sanitize_xlsx_text(value: object) -> str:
    text = str(value)
    return re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]", "", text)


def append_day_report(
    doc: Document,
    report: DayReport,
    screenshots_dir: Path | None,
    auto_result_images: bool = True,
) -> None:
    add_day_heading(doc, f"第{chinese_number(report.day_index)}天：{report.title}")

    add_section_heading(doc, "1、课程目标：")
    for index, goal in enumerate(report.goals, start=1):
        add_body_paragraph(doc, f"（{index}）{goal}")

    add_section_heading(doc, "2、课程内容：")
    add_body_paragraph(doc, "；".join(report.topics) + "。")

    add_section_heading(doc, "3、课程内容详情：")
    for index, detail in enumerate(report.details, start=1):
        add_body_paragraph(doc, f"（{index}）{detail}")

    add_section_heading(doc, "4、课程代码及执行过程：")
    if not report.tasks:
        add_body_paragraph(doc, NO_TASK_MESSAGE)
        return
    for task_index, task in enumerate(report.tasks, start=1):
        add_task_block(doc, report.day_index, task_index, task, screenshots_dir, auto_result_images)


def add_task_block(
    doc: Document,
    day_index: int,
    task_index: int,
    task: TaskItem,
    screenshots_dir: Path | None,
    auto_result_images: bool = True,
) -> None:
    add_section_heading(doc, f"任务{task_index}：{task.title}")
    add_body_paragraph(doc, f"【任务要求】{task.requirement}")
    add_body_paragraph(doc, "【代码】")
    add_code_block(doc, task.code)
    add_body_paragraph(doc, "运行结果：")

    image_path = find_screenshot(screenshots_dir, day_index, task_index) if screenshots_dir else None
    if image_path is None and auto_result_images:
        image_path = create_result_screenshot(task, day_index, task_index)
    if image_path:
        paragraph = doc.add_paragraph(style="实训报告_图题")
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = paragraph.add_run()
        run.add_picture(str(image_path), width=Inches(5.7))
    else:
        add_caption_paragraph(doc, f"（此处插入 {day_index}-{task_index} {task.caption}运行结果截图）")
    add_caption_paragraph(doc, f"图{day_index}-{task_index} {task.caption}运行结果")


def create_result_screenshot(task: TaskItem, day_index: int, task_index: int) -> Path | None:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return None

    AUTO_IMAGE_DIR.mkdir(parents=True, exist_ok=True)
    digest = hashlib.sha256(task.code.encode("utf-8", errors="ignore")).hexdigest()[:10]
    caption = safe_image_stem(task.caption or task.title)
    target = AUTO_IMAGE_DIR / f"{day_index}-{task_index}-{caption}-{digest}.png"
    lines, exit_code = build_console_output(task, day_index, task_index)
    render_pycharm_console_image(
        lines=lines,
        exit_code=exit_code,
        target=target,
        image_module=Image,
        draw_module=ImageDraw,
        font_module=ImageFont,
    )
    return target if target.exists() else None


def build_console_output(task: TaskItem, day_index: int, task_index: int) -> tuple[list[str], int]:
    executed = execute_python_code_safely(task.code, day_index, task_index)
    if executed is not None:
        output, exit_code = executed
    else:
        output, exit_code = simulate_code_output(task), 0

    lines = [PYCHARM_COMMAND_LINE]
    lines.extend(line.rstrip() for line in output.splitlines())
    if len(lines) == 1:
        lines.append("程序运行完成。")
    lines.append("")
    lines.append(f"进程已结束，退出代码为 {exit_code}")
    return lines, exit_code


def execute_python_code_safely(code: str, day_index: int, task_index: int) -> tuple[str, int] | None:
    if not is_safe_code_for_execution(code):
        return None

    input_values = build_input_values(code, seed=day_index * 100 + task_index)
    wrapped_code = build_execution_wrapper(code, input_values, seed=day_index * 100 + task_index)
    try:
        with tempfile.TemporaryDirectory(prefix="report_code_run_") as temp_dir:
            script_path = Path(temp_dir) / "lianxi.py"
            script_path.write_text(wrapped_code, encoding="utf-8")
            env = os.environ.copy()
            env["PYTHONIOENCODING"] = "utf-8"
            result = subprocess.run(
                [sys.executable, str(script_path)],
                cwd=temp_dir,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=5,
                env=env,
                check=False,
            )
    except Exception:
        return None

    output = (result.stdout or "").rstrip()
    if result.stderr:
        output = (output + "\n" if output else "") + result.stderr.strip()
    if len(output) > 5000:
        output = output[:5000].rstrip() + "\n……输出过长，后续内容已省略。"
    return output, result.returncode


def build_execution_wrapper(code: str, input_values: list[str], seed: int) -> str:
    return "\n".join(
        [
            "import builtins",
            "import random",
            f"random.seed({seed})",
            f"_report_inputs = iter({input_values!r})",
            "def _report_input(prompt=''):",
            "    print(str(prompt), end='')",
            "    try:",
            "        value = next(_report_inputs)",
            "    except StopIteration:",
            "        value = ''",
            "    print(value)",
            "    return value",
            "builtins.input = _report_input",
            code,
        ]
    )


def build_input_values(code: str, seed: int = 0) -> list[str]:
    prompts = re.findall(r"input\(\s*(?:f)?[\"']([^\"']*)[\"']?", code)
    if "班级" in code and "组别" in code and "姓名" in code and len(prompts) < 3:
        prompts = ["请输入班级：", "请输入组别：", "请输入姓名："]
    if not prompts and "input(" in code:
        prompts = ["请输入内容"]

    rng = random.Random(f"input-values-{seed}-{hashlib.md5(code.encode('utf-8')).hexdigest()[:10]}")
    values: list[str] = []
    number_inputs = [str(value) for value in rng.sample(range(12, 96), 7)]
    for index, prompt in enumerate(prompts):
        if "用户名" in prompt:
            values.append(make_sample_username(rng))
        elif "班级" in prompt:
            values.append(rng.choice(SAMPLE_CLASSES))
        elif "组别" in prompt or "组号" in prompt:
            values.append(rng.choice(SAMPLE_GROUPS))
        elif "姓名" in prompt:
            if "用户名" in prompt:
                values.append(make_sample_username(rng))
            else:
                values.append(rng.choice(SAMPLE_NAMES if rng.random() < 0.55 else SAMPLE_CHINESE_NAMES))
        elif "单价" in prompt or "价格" in prompt:
            values.append(f"{rng.choice([8.5, 12.0, 15.8, 19.9, 26.5]):.1f}")
        elif "数量" in prompt:
            values.append(str(rng.randint(2, 6)))
        elif "商品" in prompt:
            values.append(rng.choice(SAMPLE_PRODUCTS))
        elif "数字" in prompt or "猜" in prompt:
            values.append(number_inputs[min(index, len(number_inputs) - 1)])
        else:
            values.append(rng.choice(SAMPLE_PROJECTS + SAMPLE_TEXT_LINES + SAMPLE_NAMES))

    if "猜数字" in code and len(values) < 7:
        values.extend(number_inputs[len(values) :])
    return values


def make_sample_username(rng: random.Random) -> str:
    prefix = rng.choice(["lin", "chen", "zhao", "xu", "luo", "tang", "qin", "sun"])
    suffix = str(rng.randint(10, 99))
    candidate = f"{prefix}_{suffix}"
    if len(candidate) > 8:
        candidate = candidate[:8]
    return candidate.ljust(8, "7")


def is_safe_code_for_execution(code: str) -> bool:
    try:
        tree = ast.parse(code)
    except SyntaxError:
        return False

    for node in ast.walk(tree):
        if isinstance(node, (ast.Delete, ast.Global, ast.Nonlocal)):
            return False
        if isinstance(node, ast.While) and isinstance(node.test, ast.Constant) and node.test.value is True:
            return False
        if isinstance(node, ast.Import):
            if any(alias.name.split(".")[0] not in SAFE_IMPORTS for alias in node.names):
                return False
        if isinstance(node, ast.ImportFrom):
            module = (node.module or "").split(".")[0]
            if module not in SAFE_IMPORTS:
                return False
        if isinstance(node, ast.Attribute) and node.attr.startswith("__"):
            return False
        if isinstance(node, ast.Name) and node.id in DANGEROUS_NAMES:
            return False
        if isinstance(node, ast.Call):
            if not is_safe_call_node(node):
                return False
            if is_oversized_range_call(node):
                return False
    return True


def is_oversized_range_call(node: ast.Call) -> bool:
    if not isinstance(node.func, ast.Name) or node.func.id != "range":
        return False
    numeric_args = [arg.value for arg in node.args if isinstance(arg, ast.Constant) and isinstance(arg.value, int)]
    return any(abs(value) > 10000 for value in numeric_args)


def is_safe_call_node(node: ast.Call) -> bool:
    func = node.func
    if isinstance(func, ast.Name):
        if func.id in DANGEROUS_NAMES:
            return False
        if func.id == "open":
            return is_safe_open_call(node)
    if isinstance(func, ast.Attribute):
        if func.attr.startswith("__") or func.attr in DANGEROUS_NAMES:
            return False
    return True


def is_safe_open_call(node: ast.Call) -> bool:
    if not node.args or not isinstance(node.args[0], ast.Constant) or not isinstance(node.args[0].value, str):
        return False
    file_name = node.args[0].value
    path = Path(file_name)
    if path.is_absolute() or ".." in path.parts:
        return False
    if "\\" in file_name or "/" in file_name:
        return False
    return True


def simulate_code_output(task: TaskItem) -> str:
    code = task.code
    task_text = f"{task.title} {task.requirement} {code}"
    rng = random.Random(f"simulate-{hashlib.md5(code.encode('utf-8')).hexdigest()[:12]}")
    for chart_name in ("柱状图", "折线图", "饼图", "散点图", "直方图", "箱线图", "雷达图"):
        if chart_name in task_text:
            sample_count = rng.randint(4, 8)
            return f"{chart_name}已生成，样本数量： {sample_count}\n图表标题、坐标轴或分类标签已完成设置。"
    if "请输入用户名" in code and "upper()" in code:
        username = make_sample_username(rng)
        return f"请输入用户名（包含数字、字母、下划线，长度为8）：{username}\n用户名： {username.upper()}"
    if "转换为大写" in code and "转换为小写" in code:
        sample = make_sample_username(rng)
        return f"请输入一段英文或用户名：{sample}\n原始内容： {sample}\n转换为大写： {sample.upper()}\n转换为小写： {sample.lower()}"
    if "原始新闻" in code and "足球" in code:
        return "\n".join(
            [
                "原始新闻： 在昨晚进行的一场足球比赛中，主场作战的球队以3比1战胜了客队。足球运动吸引了众多球迷观看，年轻球员的配合提升了足球进攻的效率。",
                "第1个'足球'出现的位置： 9",
                "'足球'出现的次数： 3",
                "替换后的新闻： 在昨晚进行地一场足球比赛中，主场作战地球队以3比1战胜了客队。足球运动吸引了众多球迷观看，年轻球员地配合提升了足球进攻地效率。",
            ]
        )
    if "关键词首次出现位置" in code and "关键词出现次数" in code:
        return "原始文本： Python课堂练习中，Python可以处理文本、数据和图表。\n关键词首次出现位置： 0\n关键词出现次数： 2"
    if "替换前" in code and "替换后" in code:
        return "替换前： Python实训需要记录代码过程，代码运行结果也要保存。\n替换后： Python实训需要记录程序过程，程序运行结果也要保存。"
    if "strip()" in code or ".strip()" in code:
        return "原始字符串： '\\n\\n\\n我们正在进行暑假实训\\n\\n\\n'\n去除前后空白后： '我们正在进行暑假实训'"
    if "split(" in code and "join(" in code:
        return "拆分后的列表： ['姓名:小花', '年龄:20', '性别:女', '电话:13444444444']\n组合后的字符串： 姓名:小花+年龄:20+性别:女+电话:13444444444"
    if "range(1, 10)" in code and "*{i}" in code:
        return "\n".join(
            "\t".join(f"{j}*{i}={i * j}" for j in range(1, i + 1))
            for i in range(1, 10)
        )
    if "猜数字" in code or "randint" in code:
        return "\n".join(
            [
                "第1次请输入1到100之间的数字：50",
                "您猜小了",
                "剩余次数： 5",
                "第2次请输入1到100之间的数字：75",
                "您猜大了",
                "剩余次数： 4",
                "第3次请输入1到100之间的数字：63",
                "恭喜你猜对了",
                "猜测记录： [50, 75, 63]",
            ]
        )
    if "continue" in code and "break" in code:
        return "1 3 5 7 9 11 13 15 "
    if "calc_total" in code:
        product = rng.choice(SAMPLE_PRODUCTS)
        price = rng.choice([8.5, 12.0, 15.8, 19.9, 26.5])
        count = rng.randint(2, 6)
        return f"请输入商品名称：{product}\n请输入商品单价：{price:.1f}\n请输入购买数量：{count}\n{product}的总价为：{price * count:.2f}元"
    if "文件内容如下" in code:
        return "文件内容如下：\n今天练习了Python文件读写命令。\nwith open()结构可以帮助程序自动关闭文件。"
    if "类型为" in code:
        return "\n".join(
            [
                "字符类型 ： Python编程语言 ，类型为： <class 'str'>",
                "数值类型 ： 2026 ，类型为： <class 'int'>",
                "布尔类型 ： True ，类型为： <class 'bool'>",
                "列表 ： ['西安', '北京', '上海'] ，类型为： <class 'list'>",
                "字典 ： {'姓名': '小明', '年龄': 20} ，类型为： <class 'dict'>",
                "元组 ： ('暑期实训',) ，类型为： <class 'tuple'>",
                "集合 ： {1, 2, 3} ，类型为： <class 'set'>",
            ]
        )
    if "请输入班级" in code and "请输入组别" in code and "请输入姓名" in code:
        class_name = rng.choice(SAMPLE_CLASSES)
        group = rng.choice(SAMPLE_GROUPS)
        name = rng.choice(SAMPLE_NAMES if rng.random() < 0.55 else SAMPLE_CHINESE_NAMES)
        return f"请输入班级：{class_name}\n请输入组别：{group}\n请输入姓名：{name}\n我是{class_name}班的学生，我是第{group}组，我的姓名是：{name}"

    printed = extract_simple_print_output(code)
    if printed:
        return printed
    return f"任务名称：{task.title}\n任务要求：{task.requirement}\n程序运行完成。"


def extract_simple_print_output(code: str) -> str:
    lines: list[str] = []
    try:
        tree = ast.parse(code)
    except SyntaxError:
        return ""
    for node in ast.walk(tree):
        if not isinstance(node, ast.Call) or not isinstance(node.func, ast.Name) or node.func.id != "print":
            continue
        parts: list[str] = []
        for arg in node.args:
            if isinstance(arg, ast.Constant):
                parts.append(str(arg.value))
        if parts:
            lines.append(" ".join(parts))
    return "\n".join(lines)


def render_pycharm_console_image(
    lines: list[str],
    exit_code: int,
    target: Path,
    image_module,
    draw_module,
    font_module,
) -> None:
    width = 754
    height = 351
    font = load_result_font(font_module, 16)
    title_font = load_result_font(font_module, 16)
    icon_font = load_result_font(font_module, 18)
    small_font = load_result_font(font_module, 14)
    wrapped_lines = wrap_console_lines(lines, 78)
    line_height = 28
    max_lines = 8
    if len(wrapped_lines) > max_lines:
        wrapped_lines = wrapped_lines[: max_lines - 1] + ["……输出较多，后续内容已省略。"]

    image = image_module.new("RGB", (width, height), "#17191d")
    draw = draw_module.Draw(image)

    draw.rectangle((0, 0, width, 52), fill="#1f2227")
    draw.text((8, 18), "运行", fill="#f3f6fb", font=title_font)
    draw.rounded_rectangle((68, 10, 176, 43), radius=7, fill="#243452", outline="#3b63a5", width=1)
    draw.rounded_rectangle((82, 17, 105, 36), radius=5, fill="#3777c8")
    draw.text((119, 17), "lianxi", fill="#dfe8f7", font=title_font)
    draw.text((157, 16), "×", fill="#8d96a6", font=small_font)

    draw.rectangle((0, 52, width, 97), fill="#181b20")
    draw.line((0, 52, width, 52), fill="#2a2e36", width=1)
    draw.line((0, 97, width, 97), fill="#2a2e36", width=1)
    draw.rectangle((0, 97, 44, height), fill="#181b20")
    draw.line((44, 97, 44, height), fill="#2b3038", width=1)
    draw.text((14, 66), "▷", fill="#5fb66d", font=icon_font)
    draw.text((48, 67), "■", fill="#8b929f", font=small_font)
    draw.text((88, 64), "⋮", fill="#adb4c0", font=icon_font)
    for offset, icon in enumerate(["↑", "↓", "≡", "⇩", "▣", "⌫"]):
        draw.text((14, 112 + offset * 32), icon, fill="#777e8a", font=icon_font)

    y = 104
    for line in wrapped_lines:
        color = "#bfc7d5"
        if "Traceback" in line or exit_code != 0 and line == wrapped_lines[-1]:
            color = "#ff7b72"
        if line.startswith("请输入") and "：" in line:
            prefix, value = line.rsplit("：", 1)
            draw.text((61, y), prefix + "：", fill="#bfc7d5", font=font)
            prefix_width = draw.textlength(prefix + "：", font=font)
            draw.text((61 + prefix_width, y), value, fill="#65b96d", font=font)
        else:
            draw.text((61, y), line, fill=color, font=font)
        y += line_height

    target.parent.mkdir(parents=True, exist_ok=True)
    image.save(target)


def load_result_font(font_module, size: int):
    candidates = [
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simsun.ttc",
        r"C:\Windows\Fonts\consola.ttf",
        r"C:\Windows\Fonts\arial.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            try:
                return font_module.truetype(candidate, size=size)
            except Exception:
                continue
    return font_module.load_default()


def wrap_console_lines(lines: list[str], max_chars: int) -> list[str]:
    wrapped: list[str] = []
    for line in lines:
        if len(line) <= max_chars:
            wrapped.append(line)
            continue
        current = line
        while len(current) > max_chars:
            wrapped.append(current[:max_chars])
            current = "    " + current[max_chars:]
        wrapped.append(current)
    return wrapped


def safe_image_stem(text: str) -> str:
    text = cleanup_caption(text) or "运行结果"
    text = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", text)
    return text[:24]


def add_day_heading(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph(style="实训报告_日标题")
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run(text)
    set_paragraph_format(paragraph, "day")
    set_run_font(run, size=14, bold=True)


def add_section_heading(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph(style="实训报告_小节标题")
    run = paragraph.add_run(text)
    set_paragraph_format(paragraph, "section")
    set_run_font(run, size=BODY_SIZE)


def add_body_paragraph(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph(style="实训报告_正文")
    run = paragraph.add_run(text)
    set_paragraph_format(paragraph, "body")
    set_run_font(run, size=BODY_SIZE)


def add_caption_paragraph(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph(style="实训报告_图题")
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run(text)
    set_paragraph_format(paragraph, "caption")
    set_run_font(run, size=BODY_SIZE)


def add_code_block(doc: Document, code: str) -> None:
    for line in code.splitlines():
        paragraph = doc.add_paragraph(style="实训报告_代码")
        paragraph.add_run(line if line else " ")
        set_paragraph_format(paragraph, "code")
        for run in paragraph.runs:
            set_run_font(run, size=BODY_SIZE)


def set_paragraph_format(paragraph, role: str) -> None:
    fmt = paragraph.paragraph_format
    fmt.space_before = Pt(0)
    fmt.space_after = Pt(0)
    fmt.line_spacing = 1.5
    if role == "body":
        fmt.first_line_indent = BODY_FIRST_LINE_INDENT
        fmt.left_indent = Pt(0)
    elif role == "code":
        fmt.first_line_indent = Pt(0)
        fmt.left_indent = BODY_FIRST_LINE_INDENT
    else:
        fmt.first_line_indent = Pt(0)
        fmt.left_indent = Pt(0)
    if role in {"caption", "day"}:
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER


def set_run_font(run, size: float = BODY_SIZE, bold: bool | None = None) -> None:
    run.font.name = "Times New Roman"
    run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold

    r_pr = run._element.get_or_add_rPr()
    r_fonts = r_pr.rFonts
    if r_fonts is None:
        r_fonts = OxmlElement("w:rFonts")
        r_pr.insert(0, r_fonts)
    r_fonts.set(qn("w:ascii"), "Times New Roman")
    r_fonts.set(qn("w:hAnsi"), "Times New Roman")
    r_fonts.set(qn("w:eastAsia"), "宋体")
    r_fonts.set(qn("w:cs"), "Times New Roman")


def clear_paragraph(paragraph) -> None:
    for child in list(paragraph._p):
        if child.tag != qn("w:pPr"):
            paragraph._p.remove(child)


def find_screenshot(screenshots_dir: Path | None, day_index: int, task_index: int) -> Path | None:
    if screenshots_dir is None or not screenshots_dir.exists():
        return None
    prefixes = [f"{day_index}-{task_index}", f"图{day_index}-{task_index}"]
    extensions = (".png", ".jpg", ".jpeg", ".bmp")
    matches: list[Path] = []
    for path in screenshots_dir.iterdir():
        if path.suffix.lower() not in extensions:
            continue
        if any(path.stem.startswith(prefix) for prefix in prefixes):
            matches.append(path)
    return sorted(matches)[0] if matches else None


def infer_day_title(source: Path, text: str, topics: list[str]) -> str:
    if "九九乘法表" in text and "文件读写" in text:
        return "循环结构、函数与文件读写"
    if "数据类型函数" in text or "upper()" in text or "split" in text:
        return "字符串函数与文本处理"
    if "数据类型的书写" in text and "输入/输出" in text:
        return "Python编程语言基础与数据类型元素访问"
    for topic in topics:
        if topic not in {"实训规划与注意事项", "Python编程语言与业务应用"}:
            return topic
    return cleanup_heading(source.stem)


def infer_task_title(requirement: str) -> str:
    text = cleanup_requirement(requirement)
    lowered = text.lower()
    if "柱状图" in text or "条形图" in text:
        return "柱状图生成"
    if "折线图" in text or "曲线图" in text:
        return "折线图生成"
    if "饼图" in text or "圆饼图" in text:
        return "饼图生成"
    if "散点图" in text:
        return "散点图生成"
    if "直方图" in text:
        return "直方图生成"
    if "箱线图" in text or "箱型图" in text:
        return "箱线图生成"
    if "雷达图" in text:
        return "雷达图生成"
    if "用户名" in text and "大写" in text:
        return "用户名大写输出"
    if "upper" in lowered or "lower" in lowered or ("大写" in text and "小写" in text):
        return "字符串大小写转换"
    if "足球" in text:
        return "足球新闻字符串查找、统计与替换"
    if "find" in lowered or "count" in lowered or "查找" in text or "统计" in text:
        return "字符串查找与统计"
    if "replace" in lowered or "替换" in text:
        return "字符串replace()替换"
    if "去掉" in text and "空格" in text or "strip" in text.lower():
        return "字符串strip()去空"
    if "拆分" in text or "连接成字符串" in text or "join" in text.lower():
        return "字符串split()拆分与join()组合"
    if "数据类型" in text or "type" in lowered or "列表" in text or "字典" in text:
        return "数据类型与元素访问"
    if "输入" in text and "输出" in text or "input" in lowered or "print" in lowered:
        return "输入输出程序"
    if "函数" in text or "def" in lowered or "return" in lowered:
        return "函数定义与返回值"
    if "循环" in text or "for" in lowered or "while" in lowered:
        return "循环结构程序"
    if "文件" in text or "open" in lowered or "read" in lowered or "write" in lowered:
        return "文件读写程序"
    return cleanup_caption(text[:28]) or "课堂练习"


def cleanup_heading(text: str) -> str:
    text = text.strip()
    text = re.sub(r"^[（(]?\d+[）).、]\s*", "", text)
    text = text.strip("：:；;，,。 ")
    return text


def cleanup_requirement(text: str) -> str:
    text = " ".join(text.split())
    text = re.sub(r"^[（(]?\d+[）).、]\s*", "", text)
    return text.strip()


def cleanup_caption(text: str) -> str:
    text = cleanup_requirement(text)
    text = re.sub(r"[。；;：:，,]+$", "", text)
    text = re.sub(r"\s+", "", text)
    return text[:32]


def clean_lines(text: str, keep_indent: bool = False) -> list[str]:
    lines: list[str] = []
    for raw_line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        line = raw_line.rstrip()
        if not line.strip():
            continue
        lines.append(line if keep_indent else line.strip())
    return lines


def is_assignment_line(line: str) -> bool:
    return bool(re.search(r"提交|邮箱|附件|命名|最迟|报告要求", line))


def contains_all(text: str, needles: Iterable[str]) -> bool:
    return all(needle in text for needle in needles)


def unique_keep_order(items: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for item in items:
        item = item.strip()
        if not item or item in seen:
            continue
        seen.add(item)
        result.append(item)
    return result


def merge_tasks(tasks: list[TaskItem]) -> list[TaskItem]:
    result: list[TaskItem] = []
    seen: set[str] = set()
    for task in tasks:
        key = task.caption
        if key in seen:
            continue
        seen.add(key)
        result.append(task)
    return result


def chinese_number(number: int) -> str:
    if 0 <= number <= 10:
        return CN_NUMS[number]
    if number < 20:
        return "十" + CN_NUMS[number - 10]
    tens, ones = divmod(number, 10)
    return CN_NUMS[tens] + "十" + (CN_NUMS[ones] if ones else "")


if __name__ == "__main__":
    raise SystemExit(main())
